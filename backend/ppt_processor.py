import base64
import html
import json
import re
import shutil
import time
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List

from ai_connector import call_external_optimizer, should_use_external

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


STYLE_PRESETS: Dict[str, Dict[str, str]] = {
    "teaching": {
        "label": "Teaching Blue",
        "primary": "#1A73E8",
        "accent": "#7EC8E3",
        "soft": "#E8F0FE",
        "background": "#F0F2F5",
        "surface": "#FFFFFF",
        "text": "#333333",
        "muted": "#666666",
        "border": "#E0E0E0",
        "font": "Segoe UI, Tahoma, Geneva, Verdana, Microsoft YaHei, sans-serif",
    },
    "softlesson": {
        "label": "Soft Lesson",
        "primary": "#2563eb",
        "accent": "#87ceeb",
        "soft": "#e6f7ff",
        "background": "#eef4f7",
        "surface": "#fdf8f4",
        "text": "#333333",
        "muted": "#666666",
        "border": "#87ceeb",
        "font": "Comic Sans MS, Segoe Print, Bradley Hand ITC, Segoe UI, cursive",
    },
    "webacademic": {
        "label": "Academic Webpage",
        "primary": "#1A365D",
        "accent": "#2A6F97",
        "soft": "#F0F7FF",
        "background": "#F0F7FF",
        "surface": "#FFFFFF",
        "text": "#333333",
        "muted": "#666666",
        "border": "#E0E0E0",
        "font": "Arial, Microsoft YaHei, PingFang SC, sans-serif",
    },
    "clean": {
        "label": "Clean",
        "primary": "#12356b",
        "accent": "#0b74de",
        "soft": "#eef6ff",
        "background": "#f6f9fc",
        "surface": "#ffffff",
        "text": "#0b1736",
        "muted": "#5d6b82",
        "border": "#d8e4f0",
        "font": "Inter, Segoe UI, Arial, sans-serif",
    },
    "academic": {
        "label": "Academic",
        "primary": "#1f3a5f",
        "accent": "#637a9f",
        "soft": "#f2f5f9",
        "background": "#f5f7fb",
        "surface": "#ffffff",
        "text": "#172033",
        "muted": "#5b6372",
        "border": "#d9dee8",
        "font": "Georgia, Times New Roman, serif",
    },
    "instructional": {
        "label": "Instructional",
        "primary": "#004d7a",
        "accent": "#00a6d6",
        "soft": "#eafaff",
        "background": "#f4fbfd",
        "surface": "#ffffff",
        "text": "#083044",
        "muted": "#486676",
        "border": "#cfeaf3",
        "font": "Aptos, Segoe UI, Arial, sans-serif",
    },
    "minimal": {
        "label": "Minimal",
        "primary": "#111827",
        "accent": "#2563eb",
        "soft": "#f3f4f6",
        "background": "#f7f8fa",
        "surface": "#ffffff",
        "text": "#111827",
        "muted": "#6b7280",
        "border": "#e5e7eb",
        "font": "Helvetica Neue, Arial, sans-serif",
    },
    "contrast": {
        "label": "High Contrast",
        "primary": "#07111f",
        "accent": "#00d4ff",
        "soft": "#081a2b",
        "background": "#050b14",
        "surface": "#ffffff",
        "text": "#061020",
        "muted": "#31425c",
        "border": "#b8c7de",
        "font": "Segoe UI, Arial, sans-serif",
    },
    "healing": {
        "label": "Healing Hand-drawn",
        "primary": "#5e7158",
        "accent": "#df8f75",
        "soft": "#fff3df",
        "background": "#fbf6ec",
        "surface": "#fffaf1",
        "text": "#3d3328",
        "muted": "#75695c",
        "border": "#ead9bd",
        "font": "Segoe Print, Comic Sans MS, Bradley Hand ITC, Segoe UI, cursive",
    },
    "doodle": {
        "label": "Doodle Sketch",
        "primary": "#3b2925",
        "accent": "#8fc7df",
        "soft": "#f8ebd2",
        "background": "#eaf3f8",
        "surface": "#fff1d8",
        "text": "#332522",
        "muted": "#6f5c52",
        "border": "#3b2925",
        "font": "Segoe Print, Comic Sans MS, Bradley Hand ITC, Ink Free, cursive",
    },
    "swiss": {
        "label": "Swiss Grid",
        "primary": "#111111",
        "accent": "#2458ff",
        "soft": "#f2f2f2",
        "background": "#f8f8f5",
        "surface": "#ffffff",
        "text": "#111111",
        "muted": "#5f5f5f",
        "border": "#d8d8d2",
        "font": "Arial, Helvetica Neue, Helvetica, sans-serif",
    },
    "editorial": {
        "label": "Editorial",
        "primary": "#3a2618",
        "accent": "#b45432",
        "soft": "#f8efe5",
        "background": "#f4ece2",
        "surface": "#fffaf4",
        "text": "#2f261f",
        "muted": "#76685d",
        "border": "#dfcec0",
        "font": "Georgia, Times New Roman, serif",
    },
    "vivid": {
        "label": "Vivid",
        "primary": "#25135f",
        "accent": "#ff5a3d",
        "soft": "#f1edff",
        "background": "#f7f5ff",
        "surface": "#ffffff",
        "text": "#211935",
        "muted": "#655c77",
        "border": "#ddd5f3",
        "font": "Aptos, Segoe UI, Arial, sans-serif",
    },
}

DEFAULT_STYLE = "teaching"


@dataclass
class JobResult:
    job: Dict[str, Any]
    output_dir: Path
    html_path: Path


def safe_name(name: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "-", Path(name).stem).strip("-")
    return cleaned[:80] or "presentation"


def decode_upload(file_b64: str, target: Path) -> None:
    target.parent.mkdir(parents=True, exist_ok=True)
    payload = file_b64.split(",", 1)[-1]
    target.write_bytes(base64.b64decode(payload))


def normalize_generated_html(html_code: str, output_dir: Path, job_id: str) -> str:
    html_code = re.sub(r"\bUntitled\s+slide\b", "Slide", html_code, flags=re.IGNORECASE)
    assets_dir = (output_dir / "assets").resolve()
    output_dir = output_dir.resolve()
    replacements = {
        assets_dir.as_posix() + "/": "assets/",
        str(assets_dir) + "\\": "assets/",
        str(assets_dir) + "/": "assets/",
        output_dir.as_posix() + "/assets/": "assets/",
        str(output_dir) + "\\assets\\": "assets/",
        str(output_dir) + "/assets/": "assets/",
        f"/outputs/{job_id}/assets/": "assets/",
    }
    try:
        replacements[assets_dir.as_uri() + "/"] = "assets/"
        replacements[output_dir.as_uri() + "/assets/"] = "assets/"
    except ValueError:
        pass
    for source, target in replacements.items():
        html_code = html_code.replace(source, target)
        html_code = html_code.replace(html.escape(source), target)
    html_code = re.sub(
        rf"https?://(?:127\.0\.0\.1|localhost):\d+/outputs/{re.escape(job_id)}/assets/",
        "assets/",
        html_code,
        flags=re.IGNORECASE,
    )
    html_code = ensure_ai_edit_bridge(html_code, job_id)
    return html_code


def ensure_ai_edit_bridge(html_code: str, job_id: str) -> str:
    if "window.exportEditedHtml" in html_code and "toggleEdit" in html_code:
        return html_code
    if "ai-edit-bridge" in html_code:
        return html_code
    bridge = build_ai_edit_bridge(job_id)
    if re.search(r"</body\s*>", html_code, flags=re.IGNORECASE):
        return re.sub(
            r"</body\s*>",
            lambda _match: bridge + "\n</body>",
            html_code,
            count=1,
            flags=re.IGNORECASE,
        )
    return html_code + bridge


def build_ai_edit_bridge(job_id: str) -> str:
    bridge = r"""
<!-- ai-edit-bridge -->
<style id="aiEditBridgeStyle">
  #aiEditLayer {
    position: fixed;
    inset: 18px 18px auto auto;
    z-index: 2147483000;
    font-family: Inter, Segoe UI, Arial, sans-serif;
    color: #172033;
  }
  #aiEditLayer * { box-sizing: border-box; }
  .ai-edit-toggle,
  .ai-edit-toolbar button,
  .ai-edit-toolbar select,
  .ai-edit-toolbar input {
    font: 700 13px/1.1 Inter, Segoe UI, Arial, sans-serif;
  }
  .ai-edit-toggle {
    border: 1px solid rgba(37, 99, 235, 0.32);
    border-radius: 999px;
    padding: 11px 16px;
    background: rgba(255, 255, 255, 0.92);
    color: #2563eb;
    box-shadow: 0 14px 34px rgba(15, 35, 70, 0.18);
    cursor: pointer;
    backdrop-filter: blur(16px);
  }
  .ai-edit-toolbar {
    display: none;
    width: min(920px, calc(100vw - 36px));
    max-height: min(42vh, 320px);
    overflow: auto;
    gap: 8px;
    flex-wrap: wrap;
    align-items: center;
    padding: 12px;
    border: 1px solid rgba(37, 99, 235, 0.22);
    border-radius: 14px;
    background: rgba(255,255,255,0.96);
    box-shadow: 0 18px 48px rgba(15, 35, 70, 0.2);
    backdrop-filter: blur(18px);
  }
  body.editing #aiEditLayer {
    inset: 14px 14px auto 14px;
  }
  body.editing .ai-edit-toggle {
    display: none;
  }
  body.editing .ai-edit-toolbar {
    display: flex;
  }
  .ai-edit-toolbar button,
  .ai-edit-toolbar select,
  .ai-edit-toolbar input[type="number"] {
    min-height: 34px;
    border: 1px solid rgba(37, 99, 235, 0.22);
    border-radius: 10px;
    padding: 8px 10px;
    background: #fff;
    color: #1d2a44;
  }
  .ai-edit-toolbar button {
    cursor: pointer;
  }
  .ai-edit-toolbar .primary {
    background: #2563eb;
    color: #fff;
  }
  .ai-edit-toolbar .danger {
    color: #dc2626;
  }
  body :where(h1,h2,h3,h4,h5,h6,p,li,td,th,span,a,figcaption,blockquote,small,strong,em,b,i,u,label,button,div,section,article,main,aside,header,footer):not(#aiEditLayer):not(#aiEditLayer *):not([data-ai-user-font]):not([data-ai-nav-control="true"]):not([data-ai-nav-control="true"] *) {
    font-size: max(27pt, 1em) !important;
    line-height: 1.25 !important;
  }
  body :where(h1,h2,h3,.slide-title,.section-title,[data-title="true"],[data-ai-title="true"]):not(#aiEditLayer):not(#aiEditLayer *):not([data-ai-user-font]):not([data-ai-nav-control="true"]):not([data-ai-nav-control="true"] *) {
    font-size: max(46pt, 1em) !important;
    line-height: 1.12 !important;
  }
  body :where(input,textarea,select):not(#aiEditLayer input):not(#aiEditLayer textarea):not(#aiEditLayer select):not([data-ai-user-font]):not([data-ai-nav-control="true"]):not([data-ai-nav-control="true"] *) {
    font-size: max(27pt, 1em) !important;
    line-height: 1.25 !important;
  }
  body [data-ai-nav-control="true"] {
    min-width: 0 !important;
    min-height: 0 !important;
    width: auto !important;
    height: auto !important;
    max-height: 28px !important;
    padding: 3px 7px !important;
    border-radius: 999px !important;
    font: 700 8pt/1 Inter, Segoe UI, Arial, sans-serif !important;
    font-size: 8pt !important;
    letter-spacing: 0 !important;
    transform: none !important;
    z-index: 2147481200 !important;
  }
  body [data-ai-nav-control-wrap="true"] {
    z-index: 2147481100 !important;
  }
  body.export-scroll [data-ai-nav-control="true"] {
    display: none !important;
  }
  body.export-scroll {
    width: 100% !important;
    min-height: 100% !important;
    height: auto !important;
    overflow: auto !important;
    scroll-behavior: smooth;
  }
  body.export-scroll {
    display: block !important;
    position: relative !important;
  }
  body.export-scroll :where(.slide,.page,.deck-slide,.presentation-slide,.screen,section,article,[data-slide],[id^="slide-"]) {
    position: relative !important;
    inset: auto !important;
    display: flex !important;
    opacity: 1 !important;
    visibility: visible !important;
    pointer-events: auto !important;
    transform: none !important;
    z-index: auto !important;
    width: min(100%, 1280px) !important;
    max-width: 1280px !important;
    min-height: 100vh !important;
    height: auto !important;
    margin: 0 auto 32px !important;
    overflow: visible !important;
  }
  body.export-scroll :where(.slide:not(.active),.page:not(.active),.deck-slide:not(.active),.presentation-slide:not(.active)) {
    display: flex !important;
  }
  body.export-scroll :where(.controls,.navigation,.nav,.pager,.pagination,.slide-controls,.deck-controls,[data-ai-nav-control="true"]) {
    display: none !important;
  }
  .ai-color-swatch {
    width: 28px;
    min-width: 28px;
    padding: 0 !important;
    background: var(--swatch) !important;
  }
  .ai-color-swatch.white {
    box-shadow: inset 0 0 0 1px #cbd5e1;
  }
  body.editing [data-ai-editable="true"] {
    outline: 1px dashed rgba(37,99,235,0.3);
    outline-offset: 2px;
    cursor: text !important;
    pointer-events: auto !important;
    user-select: text !important;
    -webkit-user-select: text !important;
  }
  body.editing [contenteditable="true"] {
    cursor: text !important;
    pointer-events: auto !important;
    user-select: text !important;
    -webkit-user-select: text !important;
  }
  body.editing img.ai-edit-selected,
  body.editing .ai-edit-textbox.ai-edit-selected {
    outline: 3px solid #2563eb !important;
    outline-offset: 3px;
  }
  .ai-edit-textbox {
    position: absolute;
    z-index: 2147482500;
    min-width: 120px;
    min-height: 48px;
    padding: 10px 12px;
    border: 1px dashed #2563eb;
    border-radius: 10px;
    background: rgba(255,255,255,0.88);
    color: #172033;
    font: 700 27pt/1.25 Inter, Segoe UI, Arial, sans-serif;
    overflow-wrap: anywhere;
    cursor: move;
  }
</style>
<div id="aiEditLayer" aria-label="AI HTML editor">
  <button class="ai-edit-toggle" type="button" onclick="toggleEdit()">Edit HTML</button>
  <div class="ai-edit-toolbar" id="aiEditToolbar">
    <button type="button" class="primary" onclick="toggleEdit()">Stop Editing</button>
    <button type="button" onclick="execEdit('bold')"><strong>Bold</strong></button>
    <button type="button" onclick="execEdit('italic')"><em>Italic</em></button>
    <button type="button" onclick="execEdit('underline')"><u>Underline</u></button>
    <button type="button" onclick="execEdit('justifyLeft')">Left</button>
    <button type="button" onclick="execEdit('justifyCenter')">Center</button>
    <button type="button" onclick="execEdit('justifyRight')">Right</button>
    <select id="aiEditFont" onchange="setFontFamily(this.value)">
      <option value="">Font</option>
      <option value="Arial, sans-serif">Arial</option>
      <option value="Inter, Segoe UI, sans-serif">Inter</option>
      <option value="Georgia, serif">Georgia</option>
      <option value="'Times New Roman', serif">Times</option>
      <option value="'Microsoft YaHei', sans-serif">Microsoft YaHei</option>
      <option value="'Comic Sans MS', cursive">Comic Sans</option>
    </select>
    <input id="aiEditSize" type="number" min="1" value="28" title="Font size">
    <button type="button" onclick="setFontSize()">Size</button>
    <button class="ai-color-swatch" style="--swatch:#000000" onclick="setTextColor('#000000')" title="Black"></button>
    <button class="ai-color-swatch white" style="--swatch:#ffffff" onclick="setTextColor('#ffffff')" title="White"></button>
    <button class="ai-color-swatch" style="--swatch:#ef4444" onclick="setTextColor('#ef4444')" title="Red"></button>
    <button class="ai-color-swatch" style="--swatch:#f97316" onclick="setTextColor('#f97316')" title="Orange"></button>
    <button class="ai-color-swatch" style="--swatch:#facc15" onclick="setTextColor('#facc15')" title="Yellow"></button>
    <button class="ai-color-swatch" style="--swatch:#22c55e" onclick="setTextColor('#22c55e')" title="Green"></button>
    <button class="ai-color-swatch" style="--swatch:#2563eb" onclick="setTextColor('#2563eb')" title="Blue"></button>
    <button class="ai-color-swatch" style="--swatch:#7c3aed" onclick="setTextColor('#7c3aed')" title="Purple"></button>
    <input id="aiEditCustomColor" type="color" value="#2563eb" onchange="setTextColor(this.value)" title="Custom color">
    <button type="button" onclick="addTextBox()">Add Text Box</button>
    <button type="button" onclick="addImageByUrl()">Add Image</button>
    <button type="button" onclick="resizeSelectedImage(1.1)">Image +</button>
    <button type="button" onclick="resizeSelectedImage(0.9)">Image -</button>
    <button type="button" class="danger" onclick="deleteSelected()">Delete</button>
    <button type="button" onclick="saveEditedToServer(false)">Save Paged + Scroll</button>
    <button type="button" onclick="saveEditedToServer(true)">Download ZIP (Both)</button>
  </div>
</div>
<script id="aiEditBridgeScript">
(function() {
  const deckJobId = __JOB_ID__;
  let selected = null;
  let savedRange = null;
  const navTextPattern = /^(?:[‹›<>←→◀▶▲▼]?\s*)?(?:[‹›<>←→◀▶▲▼]|prev|previous|back|next|上一页?|下一页?|前一页?|后一页?|上一步|下一步|previous\s+slide|next\s+slide|previous\s+page|next\s+page)(?:\s*[‹›<>←→◀▶▲▼])?$/i;

  function markNavigationControls(root = document) {
    const controls = Array.from(root.querySelectorAll('button,a,[role="button"],.prev,.previous,.next,.pager,.pagination,.nav-button,.slide-nav,.deck-nav'));
    controls.forEach(control => {
      if (control.closest && control.closest('#aiEditLayer')) return;
      const text = (control.textContent || control.getAttribute('aria-label') || control.getAttribute('title') || '').replace(/\s+/g, ' ').trim();
      const className = String(control.className || '').toLowerCase();
      const id = String(control.id || '').toLowerCase();
      const isNavText = navTextPattern.test(text) || /^(?:←|→|‹|›)\s*(?:prev|previous|next)?$/i.test(text);
      const isNavMeta = /(^|[-_\s])(prev|previous|next)([-_\s]|$)/i.test(className + ' ' + id);
      if (isNavText || isNavMeta) {
        control.dataset.aiNavControl = 'true';
        control.setAttribute('data-ai-user-font', 'true');
        const wrapper = control.closest('.controls,.navigation,.nav,.pager,.pagination,.slide-controls,.deck-controls,.slide-nav,.deck-nav,[data-navigation]');
        if (wrapper && !wrapper.closest('#aiEditLayer')) {
          wrapper.dataset.aiNavControlWrap = 'true';
          wrapper.setAttribute('data-ai-user-font', 'true');
        }
      }
    });
  }

  function editableTargets() {
    const base = Array.from(document.body.querySelectorAll('h1,h2,h3,h4,h5,h6,p,li,td,th,figcaption,blockquote,span,a,button'));
    const leafBlocks = Array.from(document.body.querySelectorAll('main div, section div, article div, body > div')).filter(el => {
      if (el.closest('#aiEditLayer') || el.closest('script,style')) return false;
      if (!el.textContent || !el.textContent.trim()) return false;
      return !el.querySelector('h1,h2,h3,h4,h5,h6,p,li,td,th,table,img,section,article,main,div');
    });
    const titleLike = Array.from(document.body.querySelectorAll('[class*="title"],[class*="heading"],[data-title],[data-ai-title]')).filter(el => {
      if (el.closest('#aiEditLayer') || el.closest('script,style')) return false;
      return Boolean(el.textContent && el.textContent.trim());
    });
    return [...base, ...leafBlocks, ...titleLike].filter((el, index, all) => {
      if (el.closest('#aiEditLayer') || el.closest('script,style')) return false;
      return all.indexOf(el) === index;
    });
  }

  function toggleEdit(force) {
    const editing = force === undefined ? !document.body.classList.contains('editing') : Boolean(force);
    document.body.classList.toggle('editing', editing);
    editableTargets().forEach(el => {
      el.contentEditable = editing ? 'true' : 'false';
      if (editing) el.dataset.aiEditable = 'true';
      else el.removeAttribute('data-ai-editable');
    });
    markNavigationControls();
    bindImageEditing();
    bindTextBoxDragging();
    if (!editing) clearSelection();
  }

  function bindImageEditing() {
    document.querySelectorAll('img').forEach(img => {
      if (img.closest('#aiEditLayer') || img.dataset.aiImageBound) return;
      img.dataset.aiImageBound = 'true';
      img.addEventListener('click', event => {
        if (!document.body.classList.contains('editing')) return;
        event.preventDefault();
        event.stopPropagation();
        selectElement(img);
      });
    });
  }

  function selectElement(el) {
    clearSelection();
    selected = el;
    selected.classList.add('ai-edit-selected');
  }

  function clearSelection() {
    document.querySelectorAll('.ai-edit-selected').forEach(el => el.classList.remove('ai-edit-selected'));
    selected = null;
  }

  document.addEventListener('selectionchange', () => {
    if (!document.body.classList.contains('editing')) return;
    const selection = window.getSelection();
    if (!selection || !selection.rangeCount) return;
    const anchor = selection.anchorNode;
    if (anchor && anchor.parentElement?.closest('#aiEditLayer')) return;
    savedRange = selection.getRangeAt(0).cloneRange();
  });

  document.getElementById('aiEditToolbar')?.addEventListener('mousedown', event => {
    if (event.target.closest('button')) event.preventDefault();
  });

  document.addEventListener('click', event => {
    if (!document.body.classList.contains('editing')) return;
    if (event.target.closest('#aiEditLayer')) return;
    if (event.target.tagName !== 'IMG' && !event.target.closest('.ai-edit-textbox')) clearSelection();
  });

  function restoreSelection() {
    if (!savedRange) return false;
    const selection = window.getSelection();
    selection.removeAllRanges();
    selection.addRange(savedRange);
    return true;
  }

  function execEdit(command, value = null) {
    restoreSelection();
    document.execCommand(command, false, value);
  }

  function wrapSelection(styleMap) {
    if (!restoreSelection()) return;
    const selection = window.getSelection();
    if (!selection || !selection.rangeCount || selection.isCollapsed) return;
    const range = selection.getRangeAt(0);
    const span = document.createElement('span');
    Object.assign(span.style, styleMap);
    if (styleMap.fontSize) span.dataset.aiUserFont = 'true';
    span.appendChild(range.extractContents());
    range.insertNode(span);
    range.selectNodeContents(span);
    selection.removeAllRanges();
    selection.addRange(range);
    savedRange = range.cloneRange();
  }

  function setTextColor(color) {
    wrapSelection({ color });
  }

  function setFontSize() {
    const px = Math.max(1, parseInt(document.getElementById('aiEditSize').value, 10) || 28);
    wrapSelection({ fontSize: px + 'px' });
  }

  function setFontFamily(font) {
    if (font) wrapSelection({ fontFamily: font });
  }

  function resizeSelectedImage(scale) {
    if (!selected || selected.tagName !== 'IMG') return;
    const rect = selected.getBoundingClientRect();
    selected.style.width = Math.max(40, rect.width * scale) + 'px';
    selected.style.height = 'auto';
  }

  function deleteSelected() {
    if (!selected) return;
    const target = selected;
    clearSelection();
    target.remove();
  }

  function addImageByUrl() {
    const src = prompt('Image URL or relative assets/... path');
    if (!src) return;
    const img = document.createElement('img');
    img.src = src;
    img.alt = 'Added image';
    img.style.maxWidth = '320px';
    img.style.height = 'auto';
    img.style.position = 'absolute';
    img.style.left = '120px';
    img.style.top = (window.scrollY + 120) + 'px';
    img.style.zIndex = '2147482400';
    document.body.appendChild(img);
    bindImageEditing();
    selectElement(img);
  }

  function addTextBox() {
    const box = document.createElement('div');
    box.className = 'ai-edit-textbox';
    box.contentEditable = document.body.classList.contains('editing') ? 'true' : 'false';
    box.textContent = 'New text';
    box.style.left = '120px';
    box.style.top = (window.scrollY + 140) + 'px';
    document.body.appendChild(box);
    bindTextBoxDragging();
    selectElement(box);
  }

  function bindTextBoxDragging() {
    document.querySelectorAll('.ai-edit-textbox').forEach(box => {
      if (box.dataset.aiDragBound) return;
      box.dataset.aiDragBound = 'true';
      box.addEventListener('pointerdown', event => {
        if (!document.body.classList.contains('editing')) return;
        if (event.target !== box) return;
        selectElement(box);
        const startX = event.clientX;
        const startY = event.clientY;
        const startLeft = parseFloat(box.style.left || box.offsetLeft || 0);
        const startTop = parseFloat(box.style.top || box.offsetTop || 0);
        box.setPointerCapture(event.pointerId);
        const move = moveEvent => {
          box.style.left = startLeft + moveEvent.clientX - startX + 'px';
          box.style.top = startTop + moveEvent.clientY - startY + 'px';
        };
        const up = () => {
          box.removeEventListener('pointermove', move);
          box.removeEventListener('pointerup', up);
        };
        box.addEventListener('pointermove', move);
        box.addEventListener('pointerup', up);
      });
    });
  }

  function cleanClone(clone, mode) {
    const body = clone.querySelector('body');
    body?.classList.remove('editing', 'embedded');
    markNavigationControls(clone);
    clone.querySelectorAll('[contenteditable]').forEach(el => el.setAttribute('contenteditable', 'false'));
    clone.querySelectorAll('[data-ai-editable]').forEach(el => el.removeAttribute('data-ai-editable'));
    clone.querySelectorAll('.ai-edit-selected').forEach(el => el.classList.remove('ai-edit-selected'));
    clone.querySelectorAll('[data-ai-image-bound],[data-ai-drag-bound]').forEach(el => {
      el.removeAttribute('data-ai-image-bound');
      el.removeAttribute('data-ai-drag-bound');
    });
    if (mode === 'scroll') prepareScrollExport(clone);
    else body?.classList.remove('export-scroll');
  }

  function prepareScrollExport(clone) {
    const body = clone.querySelector('body');
    if (!body) return;
    body.classList.add('export-scroll');
    const style = clone.ownerDocument.createElement('style');
    style.id = 'aiScrollExportStyle';
    style.textContent = `
      body.export-scroll [data-ai-scroll-prepared="true"] {
        display: flex !important;
        opacity: 1 !important;
        visibility: visible !important;
        transform: none !important;
      }
    `;
    clone.querySelector('head')?.appendChild(style);
    clone.querySelectorAll('section,article,.slide,.page,.deck-slide,.presentation-slide,.screen,[data-slide],[id^="slide-"]').forEach(el => {
      el.dataset.aiScrollPrepared = 'true';
      el.classList.add('active');
      el.removeAttribute('hidden');
      el.removeAttribute('aria-hidden');
      el.style.display = '';
      el.style.opacity = '';
      el.style.visibility = '';
      el.style.transform = '';
    });
  }

  async function exportEditedHtml(mode = 'paged') {
    if (document.activeElement && document.activeElement.blur) document.activeElement.blur();
    const clone = document.documentElement.cloneNode(true);
    cleanClone(clone, mode);
    return '<!DOCTYPE html>\\n' + clone.outerHTML;
  }

  async function saveEditedToServer(downloadZip = false) {
    if (!deckJobId) {
      alert('This HTML cannot be saved back because it has no job id.');
      return false;
    }
    try {
      const response = await fetch('/api/jobs/' + encodeURIComponent(deckJobId) + '/save-edited', {
        method: 'POST',
        headers: { 'Content-Type': 'application/json' },
        body: JSON.stringify({
          pagedHtml: await exportEditedHtml('paged'),
          scrollHtml: await exportEditedHtml('scroll'),
        }),
      });
      const data = await response.json().catch(() => ({}));
      if (!response.ok) throw new Error(data.message || data.error || 'Could not save edited HTML');
      if (downloadZip) window.location.href = '/api/jobs/' + encodeURIComponent(deckJobId) + '/download?v=' + Date.now();
      else alert('Edited HTML has been saved.');
      return true;
    } catch (error) {
      alert(error.message || 'Could not save edited HTML');
      return false;
    }
  }

  window.toggleEdit = toggleEdit;
  window.exportEditedHtml = exportEditedHtml;
  window.saveEditedToServer = saveEditedToServer;
  window.execEdit = execEdit;
  window.setTextColor = setTextColor;
  window.setFontSize = setFontSize;
  window.setFontFamily = setFontFamily;
  window.resizeSelectedImage = resizeSelectedImage;
  window.deleteSelected = deleteSelected;
  window.addImageByUrl = addImageByUrl;
  window.addTextBox = addTextBox;
  markNavigationControls();
  bindImageEditing();
  bindTextBoxDragging();
})();
</script>
"""
    return bridge.replace("__JOB_ID__", json.dumps(job_id))


def clean_title_text(value: str) -> str:
    cleaned = re.sub(r"[\u25a0-\u25a3\u25a6-\u25ab\u25ad-\u25b1\ufffd\uf0a7\uf0b7]+", "", value)
    cleaned = re.sub(r"^[\s?？\u2022\u2043\u2219\u25cf\u25cb\u25e6\-–—:：]+", "", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    return cleaned or "Slide"


def is_page_number_marker(value: str, page: int | None = None) -> bool:
    if page is None:
        return False
    translated = str(value or "").translate(str.maketrans("０１２３４５６７８９", "0123456789"))
    text = re.sub(r"\s+", " ", translated).strip()
    if not text or len(text) > 28:
        return False
    patterns = [
        r"^[\-–—\s]*(\d{1,3})[\-–—\s]*$",
        r"^(?:p|page|slide|第)\.?\s*(\d{1,3})\s*(?:页)?$",
        r"^第\s*(\d{1,3})\s*页(?:\s*/\s*共?\s*\d{1,3}\s*页?)?$",
        r"^(\d{1,3})\s*/\s*\d{1,3}$",
    ]
    for pattern in patterns:
        match = re.match(pattern, text, flags=re.IGNORECASE)
        if match and int(match.group(1)) == page:
            return True
    return False


def process_ppt(
    filename: str,
    file_b64: str,
    style: str,
    app_root: Path,
    options: Dict[str, Any] | None = None,
) -> JobResult:
    options = options or {}
    if style not in STYLE_PRESETS:
        style = DEFAULT_STYLE
    job_id = f"JOB-{time.strftime('%Y%m%d-%H%M%S')}-{uuid.uuid4().hex[:6].upper()}"
    clean_source_name = safe_name(filename)
    data_root = Path(str(options.get("dataRoot") or (app_root / "data")))
    upload_path = data_root / "uploads" / f"{job_id}-{clean_source_name}.pptx"
    output_dir = data_root / "outputs" / job_id
    output_dir.mkdir(parents=True, exist_ok=True)
    assets_dir = output_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    decode_upload(file_b64, upload_path)

    slides = extract_slides(upload_path, assets_dir)
    integration_config = options.get("integrationConfig") or {}
    ai_status: Dict[str, Any] = {"mode": integration_config.get("mode", "local"), "used": False}
    preset = STYLE_PRESETS.get(style, STYLE_PRESETS[DEFAULT_STYLE])

    if should_use_external(integration_config):
        try:
            external_result = call_external_optimizer(
                integration_config,
                Path(filename).stem,
                slides,
                style,
                preset,
                options,
            )
            optimized_slides = merge_external_slides(slides, external_result.get("slides") or [])
            if external_result.get("html_code"):
                html_code = external_result["html_code"]
            else:
                html_code = render_html_deck(
                    title=Path(filename).stem,
                    slides=optimized_slides,
                    style_key=style,
                    options=options,
                    job_id=job_id,
                )
            ai_status = {
                "mode": integration_config.get("mode"),
                "used": True,
                "provider": external_result.get("provider"),
                "resultType": external_result.get("resultType"),
                "rawSummary": external_result.get("rawSummary"),
            }
        except Exception as exc:
            if not integration_config.get("fallbackToLocal", True):
                raise
            html_code = render_html_deck(
                title=Path(filename).stem,
                slides=slides,
                style_key=style,
                options=options,
                job_id=job_id,
            )
            ai_status = {
                "mode": integration_config.get("mode"),
                "used": False,
                "fallback": True,
                "error": str(exc),
            }
    else:
        html_code = render_html_deck(
            title=Path(filename).stem,
            slides=slides,
            style_key=style,
            options=options,
            job_id=job_id,
        )

    html_code = normalize_generated_html(html_code, output_dir, job_id)
    html_path = output_dir / "index.html"
    html_path.write_text(html_code, encoding="utf-8-sig")
    scroll_html_path = output_dir / "index-scroll.html"
    scroll_html_path.write_text(build_scroll_export_html(html_code), encoding="utf-8-sig")

    job = {
        "id": job_id,
        "fileName": filename,
        "slides": len(slides),
        "style": STYLE_PRESETS.get(style, STYLE_PRESETS[DEFAULT_STYLE])["label"],
        "styleKey": style,
        "status": "completed",
        "progress": 100,
        "updatedAt": time.strftime("%Y-%m-%d %H:%M:%S"),
        "previewUrl": f"/outputs/{job_id}/index.html",
        "scrollUrl": f"/outputs/{job_id}/index-scroll.html",
        "downloadUrl": f"/api/jobs/{job_id}/download",
        "outputPath": str(html_path),
        "constraints": {
            "keepTextUnchanged": bool(options.get("keepText", True)),
            "readable16px": bool(options.get("readableText", True)),
            "imagesIntact": bool(options.get("imagesIntact", True)),
        },
        "aiStatus": ai_status,
    }
    return JobResult(job=job, output_dir=output_dir, html_path=html_path)


def build_scroll_export_html(html_code: str) -> str:
    def add_export_scroll(match: re.Match[str]) -> str:
        tag = match.group(0)
        class_match = re.search(r'\bclass\s*=\s*(["\'])(.*?)\1', tag, flags=re.IGNORECASE)
        if class_match:
            classes = class_match.group(2)
            if "export-scroll" in classes.split():
                return tag
            updated = f'class={class_match.group(1)}{(classes + " export-scroll").strip()}{class_match.group(1)}'
            return tag[: class_match.start()] + updated + tag[class_match.end() :]
        return tag[:-1] + ' class="export-scroll">'

    scroll_html = re.sub(
        r"<body\b[^>]*>",
        add_export_scroll,
        html_code,
        count=1,
        flags=re.IGNORECASE,
    )
    scroll_html = re.sub(
        r'<button class="primary" onclick="downloadEdited\(\'paged\'\)">Download Paged HTML</button>\s*'
        r'<button onclick="downloadEdited\(\'scroll\'\)">Download Scroll HTML</button>',
        "",
        scroll_html,
        count=1,
    )
    return scroll_html


def merge_external_slides(
    original_slides: List[Dict[str, Any]],
    external_slides: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    if not external_slides:
        return original_slides
    by_page: Dict[int, Dict[str, Any]] = {}
    for index, slide in enumerate(external_slides, start=1):
        if not isinstance(slide, dict):
            continue
        try:
            page = int(slide.get("page") or index)
        except Exception:
            page = index
        by_page[page] = slide

    merged: List[Dict[str, Any]] = []
    used_pages = set()
    for index, original in enumerate(original_slides, start=1):
        original_page = int(original.get("page") or index)
        candidate = by_page.get(original_page, {})
        used_pages.add(original_page)
        next_slide = dict(original)
        texts = filter_page_number_texts(normalize_external_texts(candidate), original_page)
        if texts:
            next_slide["texts"] = texts
        images = normalize_external_images(candidate.get("images"), original.get("images", []))
        next_slide["images"] = images
        hints = normalize_layout_hints(candidate)
        if hints:
            next_slide["layoutHints"] = hints
        merged.append(next_slide)
    for page in sorted(page for page in by_page if page not in used_pages):
        candidate = by_page[page]
        texts = filter_page_number_texts(normalize_external_texts(candidate), page)
        if not texts:
            continue
        merged.append(
            {
                "page": page,
                "texts": texts,
                "images": normalize_external_images(candidate.get("images"), []),
                "layoutHints": normalize_layout_hints(candidate),
                "generatedKind": candidate.get("generatedKind") or candidate.get("kind") or "",
            }
        )
    return merged


def filter_page_number_texts(texts: List[Dict[str, Any]], page: int | None) -> List[Dict[str, Any]]:
    return [
        item
        for item in texts
        if item.get("kind") == "table" or not is_page_number_marker(str(item.get("content") or ""), page)
    ]


def normalize_external_texts(slide: Dict[str, Any]) -> List[Dict[str, Any]]:
    texts: List[Dict[str, Any]] = []
    raw_texts = slide.get("texts")
    if isinstance(raw_texts, list):
        for item in raw_texts:
            if isinstance(item, str):
                content = item.strip()
                if content:
                    texts.append({"content": content})
            elif isinstance(item, dict):
                content = str(item.get("content") or item.get("text") or "").strip()
                raw_rows = item.get("rows") or item.get("tableRows") or item.get("table")
                table_rows = normalize_table_rows(raw_rows)
                if table_rows:
                    texts.append(
                        {
                            "content": table_rows_to_text(table_rows),
                            "isTitle": False,
                            "kind": "table",
                            "rows": table_rows,
                        }
                    )
                elif content:
                    texts.append({"content": content, "isTitle": bool(item.get("isTitle"))})

    if not texts:
        title = str(slide.get("title") or "").strip()
        body = slide.get("body") or slide.get("bullets") or slide.get("points")
        if title:
            texts.append({"content": title, "isTitle": True})
        if isinstance(body, str) and body.strip():
            texts.append({"content": body.strip()})
        elif isinstance(body, list):
            for item in body:
                content = str(item.get("content") if isinstance(item, dict) else item).strip()
                if content:
                    texts.append({"content": content})

    for item in texts:
        if item.get("kind") != "table":
            item["isTitle"] = True
            break
    return texts


def normalize_external_images(
    raw_images: Any,
    original_images: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    if not isinstance(raw_images, list):
        return original_images
    allowed_src = {str(image.get("src") or "") for image in original_images}
    normalized: List[Dict[str, Any]] = []
    for item in raw_images:
        if isinstance(item, str):
            src = item.strip()
            image = {"src": src, "alt": "AI arranged image"}
        elif isinstance(item, dict):
            src = str(item.get("src") or "").strip()
            image = dict(item)
        else:
            continue
        if not src:
            continue
        if allowed_src and src not in allowed_src:
            continue
        if not (src.startswith("assets/") or src.startswith("data:image/")):
            continue
        image["src"] = src
        image["alt"] = str(image.get("alt") or "AI arranged image")
        normalized.append(image)
    return normalized or original_images


def normalize_layout_hints(slide: Dict[str, Any]) -> List[str]:
    allowed = {
        "image-left",
        "image-right",
        "gallery",
        "text-heavy",
        "image-dominant",
        "density-short",
        "density-medium",
        "density-dense",
        "cards",
        "two-column",
        "timeline",
        "agenda",
        "transition",
        "visual-focus",
    }
    raw = slide.get("layout_hints") or slide.get("layoutHints") or slide.get("layout") or slide.get("layout_classes")
    values: List[str] = []
    if isinstance(raw, str):
        values = re.split(r"[\s,]+", raw)
    elif isinstance(raw, list):
        values = [str(item) for item in raw]
    return [item for item in values if item in allowed]


def extract_slides(ppt_path: Path, assets_dir: Path) -> List[Dict[str, Any]]:
    prs = Presentation(str(ppt_path))
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    slides: List[Dict[str, Any]] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        texts: List[Dict[str, Any]] = []
        images: List[Dict[str, Any]] = []
        ordered_shapes = iter_ordered_shapes(slide.shapes)

        for shape in ordered_shapes:
            for shape_item in collect_shape_items(shape):
                content = shape_item.get("content", "")
                if not content:
                    continue
                if is_page_number_marker(content, slide_index):
                    continue
                shape_item["isTitle"] = (
                    shape_item.get("kind") != "table"
                    and len(texts) == 0
                    and len(content) <= 120
                )
                texts.append(shape_item)

            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    ext = image.ext or "png"
                    image_name = f"slide-{slide_index:03d}-image-{len(images) + 1}.{ext}"
                    image_path = assets_dir / image_name
                    image_path.write_bytes(image.blob)
                    images.append(
                        {
                            "src": f"assets/{image_name}",
                            "width": int(getattr(shape, "width", 0) or 0),
                            "height": int(getattr(shape, "height", 0) or 0),
                            "left": int(getattr(shape, "left", 0) or 0),
                            "top": int(getattr(shape, "top", 0) or 0),
                            "alt": f"Extracted image {len(images) + 1}",
                        }
                    )
                except Exception:
                    continue

        if not texts and not images:
            texts.append({"content": "Slide", "isTitle": True})

        slides.append(
            {
                "page": slide_index,
                "texts": texts,
                "images": images,
                "sourceSize": {"width": slide_width, "height": slide_height},
            }
        )

    return slides


def iter_ordered_shapes(shapes: Any) -> List[Any]:
    ordered = sorted(
        list(shapes),
        key=lambda shape: (
            int(getattr(shape, "top", 0) or 0),
            int(getattr(shape, "left", 0) or 0),
        ),
    )
    flattened: List[Any] = []
    for shape in ordered:
        flattened.append(shape)
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            flattened.extend(iter_ordered_shapes(child_shapes))
    return flattened


def normalize_table_rows(raw_rows: Any) -> List[List[str]]:
    if not isinstance(raw_rows, list):
        return []
    rows: List[List[str]] = []
    for raw_row in raw_rows:
        if isinstance(raw_row, dict):
            raw_cells = raw_row.get("cells") or raw_row.get("values") or raw_row.get("columns")
        else:
            raw_cells = raw_row
        if not isinstance(raw_cells, list):
            continue
        cells = [str(cell).strip() for cell in raw_cells if str(cell).strip()]
        if cells:
            rows.append(cells)
    if len(rows) < 1:
        return []
    max_cols = max(len(row) for row in rows)
    if max_cols < 2:
        return []
    return [row + [""] * (max_cols - len(row)) for row in rows]


def table_rows_to_text(rows: List[List[str]]) -> str:
    return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)


def collect_shape_items(shape: Any) -> List[Dict[str, Any]]:
    items: List[Dict[str, Any]] = []
    try:
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                items.append({"content": text})
    except Exception:
        pass

    try:
        if getattr(shape, "has_table", False):
            rows: List[List[str]] = []
            for row in shape.table.rows:
                cells = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip()]
                if cells:
                    rows.append(cells)
            rows = normalize_table_rows(rows)
            if rows:
                items.append({"content": table_rows_to_text(rows), "kind": "table", "rows": rows})
    except Exception:
        pass

    deduped: List[Dict[str, Any]] = []
    seen = set()
    for item in items:
        key = (item.get("kind") or "text", item.get("content") or "")
        if key in seen:
            continue
        seen.add(key)
        deduped.append(item)
    return deduped


def expand_long_text_slides(slides: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    expanded: List[Dict[str, Any]] = []
    for slide in slides:
        try:
            source_page = int(slide.get("page") or len(expanded) + 1)
        except Exception:
            source_page = len(expanded) + 1
        texts = filter_page_number_texts([dict(item) for item in slide.get("texts", [])], source_page)
        images = [dict(item) for item in slide.get("images", [])]
        if not texts:
            expanded.append(dict(slide, texts=texts, images=images))
            continue

        title_item = next((item for item in texts if item.get("kind") != "table"), texts[0])
        body_items = [item for item in texts if item is not title_item]
        normalized_body: List[Dict[str, Any]] = []
        for item in body_items:
            normalized_body.extend(split_oversized_text_item(item))

        body_length = sum(len(str(item.get("content") or "")) for item in normalized_body)
        line_count = sum(max(1, len(str(item.get("content") or "").splitlines())) for item in normalized_body)
        compact_limit = 260 if images else 320
        compact_items = 2 if images else 3
        compact_lines = 4 if images else 6
        if body_length <= compact_limit and len(normalized_body) <= compact_items and line_count <= compact_lines:
            expanded.append(dict(slide, texts=texts, images=images))
            continue

        chunks = chunk_body_items(normalized_body, max_chars=220 if images else 300, max_items=2 if images else 3)
        if not chunks:
            expanded.append(dict(slide, texts=texts, images=images))
            continue

        for chunk_index, chunk in enumerate(chunks):
            clone = dict(slide)
            clone["texts"] = [dict(title_item)] + chunk
            clone["images"] = images if chunk_index == 0 else []
            hints = list(clone.get("layoutHints") or [])
            if chunk_index > 0:
                hints = [
                    hint
                    for hint in hints
                    if hint not in {"image-left", "image-right", "image-dominant", "gallery"}
                ]
                hints.extend(["text-heavy", "density-dense"])
            clone["layoutHints"] = list(dict.fromkeys(hints))
            expanded.append(clone)

    for index, slide in enumerate(expanded, start=1):
        slide["page"] = index
    return expanded


def split_oversized_text_item(item: Dict[str, Any]) -> List[Dict[str, Any]]:
    content = str(item.get("content") or "")
    if item.get("kind") == "table" or len(content) <= 260:
        return [dict(item)]
    parts = [part.strip() for part in re.split(r"(?<=[。！？.!?;；])\s+|\n{2,}", content) if part.strip()]
    if len(parts) <= 1:
        parts = [part.strip() for part in re.split(r"\n|(?<=,)\s+", content) if part.strip()]
    chunks: List[Dict[str, Any]] = []
    current: List[str] = []
    current_len = 0
    for part in parts:
        if current and current_len + len(part) > 190:
            clone = dict(item)
            clone["content"] = "\n".join(current)
            chunks.append(clone)
            current = []
            current_len = 0
        current.append(part)
        current_len += len(part)
    if current:
        clone = dict(item)
        clone["content"] = "\n".join(current)
        chunks.append(clone)
    return chunks or [dict(item)]


def add_structural_slides(title: str, slides: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    if not slides:
        return slides
    original = [dict(slide) for slide in slides]
    agenda_titles = collect_agenda_titles(title, original)
    structured: List[Dict[str, Any]] = []
    has_agenda = any(
        slide.get("generatedKind") == "agenda" or extract_slide_title(slide) in {"目录", "Agenda", "Table of Contents"}
        for slide in original[:4]
    )

    structured.append(original[0])
    if len(original) >= 3 and not has_agenda:
        structured.append(make_agenda_slide(title, agenda_titles))

    section_interval = 5
    for index, slide in enumerate(original[1:], start=2):
        if index > 2 and (index - 2) % section_interval == 0:
            structured.append(make_transition_slide(extract_slide_title(slide) or f"Section {len(structured) + 1}"))
        structured.append(slide)

    for page, slide in enumerate(structured, start=1):
        slide["page"] = page
    return structured


def collect_agenda_titles(title: str, slides: List[Dict[str, Any]]) -> List[str]:
    titles: List[str] = []
    for slide in slides[1:]:
        candidate = extract_slide_title(slide)
        if not candidate or candidate == title or candidate in titles:
            continue
        titles.append(candidate)
        if len(titles) >= 6:
            break
    if titles:
        return titles
    return ["Core Content", "Key Analysis", "Practice", "Summary"]


def extract_slide_title(slide: Dict[str, Any]) -> str:
    texts = slide.get("texts") or []
    title_item = next((item for item in texts if item.get("kind") != "table"), None)
    return clean_title_text(str(title_item.get("content") or "")) if title_item else ""


def make_agenda_slide(title: str, items: List[str]) -> Dict[str, Any]:
    return {
        "page": 0,
        "generatedKind": "agenda",
        "layoutHints": ["agenda", "cards", "two-column", "density-medium", "visual-focus"],
        "texts": [
            {"content": "Agenda", "isTitle": True},
            {"content": "\n".join(items[:6]), "isTitle": False},
        ],
        "images": [],
        "sourceTitle": title,
    }


def make_transition_slide(title: str) -> Dict[str, Any]:
    return {
        "page": 0,
        "generatedKind": "transition",
        "layoutHints": ["transition", "density-short", "visual-focus"],
        "texts": [
            {"content": title, "isTitle": True},
        ],
        "images": [],
    }


def chunk_body_items(items: List[Dict[str, Any]], max_chars: int, max_items: int) -> List[List[Dict[str, Any]]]:
    chunks: List[List[Dict[str, Any]]] = []
    current: List[Dict[str, Any]] = []
    current_chars = 0
    for item in items:
        size = len(str(item.get("content") or "")) + (140 if item.get("kind") == "table" else 0)
        if current and (current_chars + size > max_chars or len(current) >= max_items):
            chunks.append(current)
            current = []
            current_chars = 0
        current.append(dict(item))
        current_chars += size
    if current:
        chunks.append(current)
    return chunks


def render_html_deck(
    title: str,
    slides: List[Dict[str, Any]],
    style_key: str,
    options: Dict[str, Any],
    job_id: str = "",
) -> str:
    if style_key not in STYLE_PRESETS:
        style_key = DEFAULT_STYLE
    preset = STYLE_PRESETS.get(style_key, STYLE_PRESETS[DEFAULT_STYLE])
    slides = add_structural_slides(title, expand_long_text_slides(slides))
    slide_markup = "\n<!-- --- -->\n".join(render_slide(slide, preset) for slide in slides)
    slide_count = len(slides)
    data_blob = html.escape(json.dumps({"title": title, "slides": slide_count, "jobId": job_id}, ensure_ascii=False))

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{html.escape(title)} - Optimized HTML</title>
  <style>
    :root {{
      --primary: {preset["primary"]};
      --accent: {preset["accent"]};
      --soft: {preset["soft"]};
      --bg: {preset["background"]};
      --surface: {preset["surface"]};
      --text: {preset["text"]};
      --muted: {preset["muted"]};
      --border: {preset["border"]};
      --font: {preset["font"]};
    }}
    * {{ box-sizing: border-box; }}
    html, body {{ margin: 0; min-height: 100%; }}
    body {{
      font-family: var(--font);
      background: var(--bg);
      color: var(--text);
      overflow-x: hidden;
    }}
    .slide *,
    .free-text-box {{
      max-width: 100%;
      overflow-wrap: anywhere;
      word-break: break-word;
      min-width: 0;
    }}
    body.theme-webacademic {{
      background: var(--bg);
      font-family: var(--font);
    }}
    body.theme-webacademic .slide {{
      background: var(--surface);
      border: 0;
      border-radius: 8px;
      box-shadow: 0 2px 8px rgba(0,0,0,0.08);
      overflow: hidden;
    }}
    body.theme-webacademic .slide-inner {{
      padding: 60px 80px 50px;
      align-items: center;
    }}
    body.theme-webacademic .slide-title {{
      color: var(--primary);
      text-align: center;
      border-bottom: 2px solid var(--accent);
      padding-bottom: 10px;
      margin-left: auto;
      margin-right: auto;
      max-width: 100%;
    }}
    body.theme-webacademic .slide.cover .slide-inner {{
      display: flex;
      flex-direction: column;
      justify-content: center;
      align-items: center;
      text-align: center;
      background: var(--surface);
    }}
    body.theme-webacademic .slide.cover .slide-title {{
      border-bottom: 0;
      font-size: clamp(40px, 4.2vw, 58px);
    }}
    body.theme-webacademic .slide-body {{
      font-size: clamp(24px, 1.75vw, 30px);
      line-height: 1.6;
    }}
    body.theme-webacademic .slide-body p {{
      margin-bottom: 16px;
    }}
    body.theme-webacademic .accent {{
      color: #F79256;
      font-weight: 800;
    }}
    body.theme-webacademic .content-list {{
      list-style-type: disc;
      background: transparent;
      border: 0;
      box-shadow: none;
      padding-left: 42px;
    }}
    body.theme-webacademic .media-frame img {{
      max-width: 600px;
      max-height: 350px;
      object-fit: contain;
      border-radius: 6px;
      border: 1px solid var(--border);
      box-shadow: none;
      padding: 0;
      background: transparent;
    }}
    body.theme-webacademic .page-number {{
      right: 30px;
      bottom: 20px;
      font-size: 20px;
      color: var(--muted);
    }}
    body.theme-teaching {{
      background: var(--bg);
      font-family: var(--font);
    }}
    body.theme-teaching .slide {{
      background: var(--surface);
      border: 0;
      border-radius: 8px;
      box-shadow: 0 2px 8px rgba(0,0,0,0.08);
      overflow: hidden;
    }}
    body.theme-teaching .slide-inner {{
      padding: 40px 56px;
      align-items: center;
    }}
    body.theme-teaching .slide-title {{
      color: var(--primary);
      border-left: 0;
      border-bottom: 3px solid color-mix(in srgb, var(--primary), #ffffff 28%);
      padding: 0 0 12px;
      width: 100%;
      max-width: 100%;
      margin: 0 0 28px;
      text-align: left;
      font-size: clamp(38px, 3.8vw, 48px);
      line-height: 1.14;
    }}
    body.theme-teaching .slide.text-only .slide-inner {{
      align-items: start;
      padding: 56px 72px 58px;
    }}
    body.theme-teaching .slide.text-only .text-column {{
      width: 100%;
      max-width: none;
      height: 100%;
    }}
    body.theme-teaching .slide.cover .slide-inner,
    body.theme-teaching .slide.transition .slide-inner {{
      display: flex;
      flex-direction: column;
      justify-content: center;
      align-items: center;
      text-align: center;
      background:
        radial-gradient(circle at 82% 22%, rgba(126, 200, 227, 0.16), transparent 24%),
        linear-gradient(180deg, #fff, #fbfcff);
    }}
    body.theme-teaching .slide.cover .slide-title,
    body.theme-teaching .slide.transition .slide-title {{
      border-left: 0;
      padding-left: 0;
      border-bottom: 4px solid var(--accent);
      padding-bottom: 12px;
      font-size: clamp(40px, 4.4vw, 58px);
      line-height: 1.08;
      max-height: none;
      overflow: visible;
    }}
    body.theme-teaching .slide.title-long .slide-title {{
      font-size: clamp(30px, 3vw, 40px);
      line-height: 1.08;
      max-height: none;
      overflow: visible;
    }}
    body.theme-teaching .slide.title-very-long .slide-title {{
      font-size: clamp(26px, 2.65vw, 34px);
      line-height: 1.06;
      max-height: none;
      overflow: visible;
    }}
    body.theme-teaching .slide.agenda .slide-inner {{
      grid-template-columns: minmax(0, 1fr);
      padding: 56px 72px;
    }}
    body.theme-teaching .slide-body {{
      font-size: clamp(24px, 1.75vw, 30px);
      line-height: 1.5;
      gap: 16px;
      min-height: 0;
      width: 100%;
    }}
    body.theme-teaching .slide.text-only.cards .slide-body {{
      grid-template-columns: repeat(auto-fit, minmax(min(520px, 100%), 1fr));
      align-content: start;
    }}
    body.theme-teaching .slide.text-only.cards .slide-body > .body-card:only-child {{
      grid-column: 1 / -1;
    }}
    body.theme-teaching .slide.cover .slide-body,
    body.theme-teaching .slide.transition .slide-body {{
      display: flex;
      justify-content: center;
      max-width: 880px;
      text-align: center;
    }}
    body.theme-teaching .body-card {{
      background: #F5F7FA;
      border: 1px solid var(--border);
      border-radius: 8px;
      padding: 20px 24px;
      min-width: 0;
      max-height: none;
      overflow: visible;
      box-shadow: 0 8px 22px rgba(26, 115, 232, 0.06);
    }}
    body.theme-teaching .slide.text-only .body-card {{
      width: 100%;
    }}
    body.theme-teaching .slide.density-dense .body-card {{
      padding: 12px 14px;
    }}
    body.theme-teaching .body-card p {{
      display: grid;
      grid-template-columns: auto minmax(0, 1fr);
      gap: 10px;
      align-items: start;
    }}
    body.theme-teaching .content-list {{
      list-style: none;
      padding-left: 0;
    }}
    body.theme-teaching .emoji-list li,
    body.theme-teaching .guided-list li {{
      display: grid;
      grid-template-columns: auto minmax(0, 1fr);
      gap: 10px;
      margin: 8px 0;
      min-width: 0;
    }}
    body.theme-teaching .point-icon {{
      display: inline-block;
      width: 6px;
      height: 30px;
      margin-top: 2px;
      border-radius: 999px;
      background: var(--accent);
    }}
    body.theme-teaching .timeline {{
      display: grid;
      grid-template-columns: repeat(3, minmax(0, 1fr));
      gap: 14px;
    }}
    body.theme-teaching .timeline-step {{
      position: relative;
      min-height: 136px;
      padding: 22px 18px 18px;
      background: #F5F7FA;
      border: 1px solid var(--border);
      border-radius: 8px;
      font-weight: 650;
      overflow: hidden;
    }}
    body.theme-teaching .step-number {{
      display: inline-grid;
      place-items: center;
      width: 42px;
      height: 42px;
      margin-bottom: 14px;
      border-radius: 999px;
      background: var(--primary);
      color: #fff;
      font-weight: 850;
    }}
    body.theme-teaching .chapter-label {{
      display: inline-block;
      color: var(--primary);
      font-weight: 900;
      font-size: clamp(26px, 2.4vw, 32px);
      line-height: 1.15;
      letter-spacing: 0.01em;
    }}
    body.theme-teaching .slide.chapter .slide-inner {{
      align-items: center;
    }}
    body.theme-teaching .slide.chapter .slide-title {{
      border-bottom: 0;
      text-align: center;
      font-size: clamp(26px, 2.4vw, 32px);
      line-height: 1.18;
      color: var(--primary);
    }}
    body.theme-teaching .slide.chapter .slide-body {{
      justify-items: center;
      text-align: center;
    }}
    body.theme-teaching .table-scroll {{
      border: 1px solid var(--border);
      border-radius: 8px;
      box-shadow: none;
      background: #fff;
      max-height: 360px;
      overflow: auto;
    }}
    body.theme-teaching .data-table th {{
      background: var(--primary);
      color: #fff;
    }}
    body.theme-teaching .media-frame img {{
      max-width: 600px;
      max-height: 390px;
      object-fit: contain;
      border: 2px solid var(--primary);
      border-radius: 8px;
      background: #fff;
      padding: 0;
      box-shadow: 0 12px 28px rgba(26, 115, 232, 0.12);
    }}
    body.theme-teaching .visual-anchor {{
      position: absolute;
      right: 68px;
      bottom: 56px;
      width: 88px;
      height: 88px;
      border: 16px solid var(--soft);
      border-radius: 999px;
      box-shadow: inset 0 0 0 2px rgba(26, 115, 232, 0.18);
      opacity: 0.72;
      pointer-events: none;
      z-index: 0;
    }}
    body.theme-teaching .slide.has-media .visual-anchor {{
      display: none;
    }}
    body.theme-teaching .slide.cover .visual-anchor,
    body.theme-teaching .slide.transition .visual-anchor {{
      position: static;
      order: -1;
      width: 124px;
      height: 124px;
      margin-bottom: 16px;
      opacity: 0.78;
    }}
    body.theme-teaching .page-number {{
      right: 30px;
      bottom: 20px;
      font-size: 20px;
      color: var(--muted);
    }}
    body.theme-healing {{
      background:
        radial-gradient(circle at 14% 18%, rgba(223, 143, 117, 0.16), transparent 24%),
        radial-gradient(circle at 88% 12%, rgba(94, 113, 88, 0.12), transparent 26%),
        var(--bg);
    }}
    body.theme-softlesson {{
      background:
        radial-gradient(circle at 14% 12%, rgba(135, 206, 235, 0.22), transparent 24%),
        radial-gradient(circle at 86% 88%, rgba(255, 228, 181, 0.35), transparent 28%),
        var(--bg);
    }}
    body.theme-softlesson .slide {{
      background:
        linear-gradient(180deg, rgba(255,255,255,0.42), rgba(255,255,255,0.08)),
        var(--surface);
      border: 2px dashed var(--border);
      border-radius: 15px;
      box-shadow: 0 4px 15px rgba(135, 206, 235, 0.22), 0 18px 42px rgba(67, 94, 120, 0.08);
      overflow: hidden;
    }}
    body.theme-softlesson .slide::before,
    body.theme-softlesson .slide::after {{
      position: absolute;
      pointer-events: none;
      color: var(--accent);
      opacity: 0.72;
      z-index: 0;
    }}
    body.theme-softlesson .slide::before {{
      content: "✎";
      left: 24px;
      top: 18px;
      font-size: 28px;
      transform: rotate(-10deg);
    }}
    body.theme-softlesson .slide::after {{
      content: "✦";
      right: 28px;
      bottom: 28px;
      font-size: 30px;
      transform: rotate(10deg);
    }}
    body.theme-softlesson .slide.cover .slide-inner {{
      justify-items: center;
      text-align: center;
      background:
        radial-gradient(circle at 50% 30%, rgba(230, 247, 255, 0.82), transparent 32%),
        linear-gradient(180deg, rgba(255,255,255,0.32), transparent);
    }}
    body.theme-softlesson .slide-title {{
      color: var(--primary);
      text-align: center;
      text-decoration: underline wavy var(--accent) 2px;
      text-underline-offset: 9px;
      font-weight: 820;
    }}
    body.theme-softlesson .slide.cover .slide-title {{
      font-size: clamp(40px, 5vw, 64px);
      max-width: 980px;
    }}
    body.theme-softlesson .slide-body {{
      gap: 14px;
    }}
    body.theme-softlesson .slide-body p,
    body.theme-softlesson .content-list {{
      background: var(--soft);
      border: 2px dashed var(--border);
      border-radius: 10px;
      padding: 16px 20px;
      box-shadow: 0 8px 18px rgba(135, 206, 235, 0.14);
    }}
    body.theme-softlesson .content-list {{
      margin: 0;
      padding-left: 48px;
      list-style-type: "✨ ";
    }}
    body.theme-softlesson .content-list li {{
      margin: 7px 0;
    }}
    body.theme-softlesson .accent {{
      color: var(--primary);
      font-weight: 820;
    }}
    body.theme-softlesson .media-frame img {{
      border: 2px dashed var(--border);
      border-radius: 10px;
      box-shadow: 0 3px 10px rgba(0, 0, 0, 0.10);
      background: rgba(255,255,255,0.68);
      padding: 4px;
    }}
    body.theme-softlesson .table-scroll {{
      border: 2px dashed var(--border);
      background: var(--soft);
    }}
    body.theme-softlesson .data-table th {{
      background: rgba(255, 255, 255, 0.62);
      color: var(--primary);
    }}
    body.theme-healing .slide {{
      border: 2px solid var(--border);
      border-radius: 18px;
      box-shadow: 12px 14px 0 rgba(94, 113, 88, 0.12), 0 18px 44px rgba(72, 58, 42, 0.12);
    }}
    body.theme-healing .media-frame {{
      border: 2px solid var(--border);
      border-radius: 16px;
      background:
        linear-gradient(135deg, rgba(255, 255, 255, 0.72), rgba(255, 243, 223, 0.92)),
        var(--soft);
    }}
    body.theme-healing .slide-title {{
      font-weight: 760;
    }}
    body.theme-healing .slide-inner::before {{
      height: 8px;
      border-radius: 24px 18px 28px 16px;
      transform: rotate(-1deg);
    }}
    body.theme-doodle {{
      background:
        radial-gradient(circle at 12% 8%, rgba(143, 199, 223, 0.32), transparent 22%),
        radial-gradient(circle at 88% 86%, rgba(143, 199, 223, 0.26), transparent 24%),
        var(--bg);
    }}
    body.theme-doodle .slide {{
      background:
        radial-gradient(circle at 8% 14%, rgba(143, 199, 223, 0.34) 0 20px, transparent 21px),
        radial-gradient(circle at 92% 18%, rgba(143, 199, 223, 0.3) 0 26px, transparent 27px),
        radial-gradient(circle at 88% 86%, rgba(143, 199, 223, 0.28) 0 32px, transparent 33px),
        var(--surface);
      border: 10px solid var(--border);
      border-radius: 3px 9px 4px 7px;
      box-shadow: 0 24px 48px rgba(59, 41, 37, 0.18);
    }}
    body.theme-doodle .slide::before,
    body.theme-doodle .slide::after {{
      content: "";
      position: absolute;
      pointer-events: none;
      z-index: 0;
    }}
    body.theme-doodle .slide::before {{
      left: 7%;
      top: 9%;
      width: 42px;
      height: 42px;
      background:
        linear-gradient(90deg, transparent 46%, var(--accent) 47% 53%, transparent 54%),
        linear-gradient(0deg, transparent 46%, var(--accent) 47% 53%, transparent 54%);
      transform: rotate(18deg);
      opacity: 0.9;
    }}
    body.theme-doodle .slide::after {{
      right: 8%;
      bottom: 12%;
      width: 48px;
      height: 48px;
      border: 5px solid var(--border);
      border-left-color: transparent;
      border-bottom-color: transparent;
      border-radius: 999px;
      transform: rotate(-24deg);
      opacity: 0.7;
    }}
    body.theme-doodle .slide-title {{
      display: inline;
      line-height: 1.04;
      text-decoration: underline;
      text-decoration-color: var(--accent);
      text-decoration-thickness: 12px;
      text-underline-offset: -4px;
      text-decoration-skip-ink: none;
      letter-spacing: 0.01em;
    }}
    body.theme-doodle .slide-body {{
      font-weight: 650;
    }}
    body.theme-doodle .slide-inner::before {{
      height: 18px;
      width: 160px;
      bottom: 58px;
      background: var(--accent);
      border-radius: 2px 7px 3px 9px;
      transform: rotate(-1.5deg);
      opacity: 0.88;
    }}
    body.theme-doodle .media-frame img {{
      filter: drop-shadow(8px 10px 0 rgba(59, 41, 37, 0.08));
    }}
    body.theme-swiss .slide {{
      box-shadow: none;
      border: 2px solid var(--primary);
    }}
    body.theme-swiss .slide-inner {{
      background-image:
        linear-gradient(var(--border) 1px, transparent 1px),
        linear-gradient(90deg, var(--border) 1px, transparent 1px);
      background-size: 64px 64px;
    }}
    body.theme-editorial .slide-title {{
      font-weight: 650;
    }}
    body.theme-vivid .slide {{
      box-shadow: 0 26px 70px rgba(37, 19, 95, 0.18);
    }}
    .deck-shell {{
      min-height: 100vh;
      padding: 22px;
    }}
    .deck-toolbar {{
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 16px;
      max-width: 1320px;
      margin: 0 auto 18px;
      color: var(--text);
    }}
    .deck-title {{
      font-size: 18px;
      font-weight: 750;
      letter-spacing: 0;
    }}
    .deck-actions {{
      display: flex;
      align-items: center;
      gap: 8px;
    }}
    .editor-toolbar {{
      display: none;
      align-items: center;
      flex-wrap: wrap;
      gap: 8px;
      max-width: 1320px;
      margin: 0 auto 18px;
      padding: 10px;
      background: #fff;
      border: 1px solid var(--border);
      border-radius: 8px;
      box-shadow: 0 10px 28px rgba(15, 35, 70, 0.08);
    }}
    body.editing .editor-toolbar {{ display: flex; }}
    .editor-toolbar label {{
      display: inline-flex;
      align-items: center;
      gap: 6px;
      color: var(--muted);
      font-size: 13px;
      font-weight: 650;
    }}
    .editor-toolbar input,
    .editor-toolbar select {{
      height: 34px;
      border: 1px solid var(--border);
      border-radius: 7px;
      padding: 0 8px;
      background: #fff;
      color: var(--text);
      font: inherit;
    }}
    .editor-toolbar input[type="color"] {{
      width: 42px;
      padding: 3px;
    }}
    .color-palette {{
      display: inline-flex;
      align-items: center;
      gap: 5px;
    }}
    .color-swatch {{
      width: 26px;
      height: 26px;
      min-width: 26px;
      border-radius: 999px;
      border: 1px solid var(--border);
      background: var(--swatch);
      cursor: pointer;
      padding: 0;
    }}
    .color-swatch.white {{
      box-shadow: inset 0 0 0 1px #cbd5e1;
    }}
    .editor-toolbar input[type="file"] {{ display: none; }}
    .editor-toolbar .toolbar-divider {{
      width: 1px;
      height: 28px;
      background: var(--border);
      margin: 0 2px;
    }}
    .editor-toolbar .edit-hint {{
      color: var(--muted);
      font-size: 12px;
      font-weight: 650;
    }}
    button {{
      border: 1px solid var(--border);
      background: var(--surface);
      color: var(--text);
      border-radius: 8px;
      padding: 10px 13px;
      font-size: 14px;
      font-weight: 650;
      cursor: pointer;
    }}
    button.primary {{
      background: var(--accent);
      color: #fff;
      border-color: var(--accent);
    }}
    button:disabled {{
      opacity: 0.45;
      cursor: not-allowed;
    }}
    .counter {{ color: var(--muted); font-size: 14px; min-width: 82px; text-align: center; }}
    .slide-stage {{
      width: min(1280px, calc(100vw - 44px));
      margin: 0 auto;
    }}
    .slide {{
      display: none;
      width: 100%;
      aspect-ratio: 16 / 9;
      background: var(--surface);
      border: 1px solid var(--border);
      box-shadow: 0 24px 70px rgba(15, 35, 70, 0.14);
      overflow: hidden;
      position: relative;
    }}
    .slide.active {{ display: grid; }}
    .slide-inner {{
      display: grid;
      grid-template-columns: minmax(0, 1.05fr) minmax(280px, 0.95fr);
      gap: 34px;
      width: 100%;
      height: 100%;
      padding: 56px 64px;
      align-items: center;
      position: relative;
      overflow: hidden;
    }}
    .slide-inner > * {{
      min-width: 0;
      min-height: 0;
    }}
    .slide.text-only .slide-inner {{
      grid-template-columns: minmax(0, 1fr);
      align-items: start;
      padding-left: 72px;
      padding-right: 72px;
    }}
    .slide.text-only .text-column {{
      width: 100%;
      max-width: none;
      height: 100%;
    }}
    .slide.text-only.cards .slide-body {{
      width: 100%;
      grid-template-columns: repeat(auto-fit, minmax(min(500px, 100%), 1fr));
      align-content: start;
      justify-content: stretch;
    }}
    .slide.text-only.cards .slide-body > .body-card:only-child {{
      grid-column: 1 / -1;
      width: 100%;
    }}
    .slide.text-only .body-card {{
      width: 100%;
    }}
    .slide.cover .slide-inner {{
      align-items: center;
      background:
        linear-gradient(120deg, rgba(11, 116, 222, 0.08), transparent 44%),
        linear-gradient(180deg, #fff, var(--surface));
    }}
    .slide.cover .slide-title {{
      font-size: clamp(38px, 4.5vw, 58px);
      line-height: 1.08;
      max-width: min(1040px, 100%);
      max-height: none;
      overflow: visible;
    }}
    .slide.title-long .slide-title {{
      font-size: clamp(30px, 3vw, 42px);
      line-height: 1.08;
      max-height: none;
      overflow: visible;
    }}
    .slide.title-very-long .slide-title {{
      font-size: clamp(26px, 2.55vw, 36px);
      line-height: 1.06;
      max-height: none;
      overflow: visible;
    }}
    .slide.cover .slide-body {{
      margin-top: 24px;
      font-size: clamp(22px, 2vw, 30px);
      color: var(--muted);
      max-width: 780px;
    }}
    .slide.image-left .media-column {{
      order: -1;
    }}
    .slide.image-left .slide-inner {{
      grid-template-columns: minmax(280px, 0.9fr) minmax(0, 1.1fr);
    }}
    .slide.image-right .slide-inner {{
      grid-template-columns: minmax(0, 1.1fr) minmax(280px, 0.9fr);
    }}
    .slide.image-dominant.image-left .slide-inner {{
      grid-template-columns: minmax(360px, 1.22fr) minmax(0, 0.78fr);
    }}
    .slide.image-dominant.image-right .slide-inner {{
      grid-template-columns: minmax(0, 0.78fr) minmax(360px, 1.22fr);
    }}
    .slide.text-heavy .slide-inner {{
      align-items: start;
      grid-template-columns: minmax(0, 1.3fr) minmax(240px, 0.7fr);
    }}
    .slide.text-heavy.image-left .slide-inner {{
      grid-template-columns: minmax(240px, 0.72fr) minmax(0, 1.28fr);
    }}
    .slide.gallery .slide-inner {{
      align-items: start;
    }}
    .slide.density-short .slide-body {{
      font-size: clamp(24px, 2.1vw, 30px);
      line-height: 1.38;
    }}
    .slide.density-dense .slide-inner {{
      gap: 24px;
      padding-top: 42px;
      padding-bottom: 48px;
    }}
    .slide.density-dense .slide-title {{
      font-size: clamp(28px, 3.2vw, 42px);
      margin-bottom: 14px;
    }}
    .slide.density-dense .slide-body {{
      gap: 9px;
      font-size: clamp(24px, 1.45vw, 26px);
      line-height: 1.32;
    }}
    .slide.text-only.density-dense .slide-body {{
      grid-template-columns: repeat(auto-fit, minmax(min(460px, 100%), 1fr));
      align-content: start;
      column-gap: 18px;
      row-gap: 12px;
      width: 100%;
    }}
    .slide.text-only.density-dense .slide-body > * {{
      min-width: 0;
    }}
    .slide-inner::before {{
      content: "";
      position: absolute;
      left: 64px;
      bottom: 46px;
      width: 78px;
      height: 5px;
      background: var(--accent);
      border-radius: 999px;
    }}
    .text-column {{
      min-width: 0;
      max-width: 100%;
      max-height: 100%;
      overflow: visible;
      z-index: 2;
      display: grid;
      grid-template-rows: auto minmax(0, 1fr);
    }}
    .slide-title {{
      font-size: clamp(28px, 3.4vw, 50px);
      line-height: 1.12;
      margin: 0 0 22px;
      color: var(--primary);
      font-weight: 820;
      text-align: center;
      letter-spacing: 0;
      overflow-wrap: anywhere;
      max-height: 3.6em;
      overflow: hidden;
    }}
    .slide-body {{
      display: grid;
      gap: 12px;
      font-size: clamp(24px, 1.9vw, 30px);
      line-height: 1.43;
      color: var(--text);
      min-height: 0;
      max-height: 100%;
      overflow: visible;
    }}
    .slide-body p {{
      margin: 0;
      white-space: pre-line;
      overflow-wrap: anywhere;
      word-break: break-word;
      max-height: 100%;
      overflow: visible;
    }}
    .body-card {{
      min-width: 0;
      min-height: 0;
      max-height: none;
      overflow: visible;
    }}
    .body-card p {{
      display: flex;
      gap: 8px;
      align-items: flex-start;
      min-width: 0;
    }}
    .point-icon {{
      flex: 0 0 auto;
    }}
    .content-list {{
      margin: 0;
      padding-left: 1.35em;
      overflow-wrap: anywhere;
      word-break: break-word;
      max-height: none;
      overflow: visible;
    }}
    .content-list li {{
      margin: 0.28em 0;
      min-width: 0;
    }}
    .accent {{
      color: var(--accent);
      font-weight: 800;
    }}
    .chapter-label {{
      display: inline-block;
      color: var(--primary);
      font-size: 28px;
      line-height: 1.18;
      font-weight: 850;
      letter-spacing: 0.01em;
    }}
    .slide.chapter .slide-title {{
      font-size: clamp(26px, 2.4vw, 32px);
      line-height: 1.18;
    }}
    .timeline {{
      display: grid;
      gap: 12px;
    }}
    .timeline-step {{
      min-width: 0;
      overflow-wrap: anywhere;
    }}
    .table-scroll {{
      width: 100%;
      overflow-x: auto;
      margin: 8px 0;
      border-radius: 12px;
      border: 1px solid color-mix(in srgb, var(--border), var(--accent) 18%);
      background: rgba(255,255,255,0.72);
      box-shadow: 0 12px 24px rgba(15, 35, 70, 0.08);
    }}
    .data-table {{
      width: 100%;
      table-layout: fixed;
      border-collapse: collapse;
      font-size: clamp(24px, 1.55vw, 28px);
      line-height: 1.35;
      color: var(--text);
    }}
    .data-table th,
    .data-table td {{
      padding: 10px 12px;
      border-bottom: 1px solid var(--border);
      text-align: left;
      vertical-align: top;
      min-width: 88px;
      overflow-wrap: anywhere;
      word-break: break-word;
    }}
    .data-table th {{
      color: var(--primary);
      background: var(--soft);
      font-weight: 800;
    }}
    .data-table tbody tr:nth-child(even) td {{
      background: rgba(255,255,255,0.54);
    }}
    .data-table tbody tr:last-child td {{
      border-bottom: 0;
    }}
    body.theme-doodle .table-scroll {{
      border: 4px solid var(--border);
      border-radius: 3px 9px 4px 7px;
      background: rgba(255, 241, 216, 0.8);
      box-shadow: 8px 9px 0 rgba(59, 41, 37, 0.12);
    }}
    body.theme-doodle .data-table th {{
      background: rgba(143, 199, 223, 0.58);
    }}
    .media-column {{
      display: grid;
      gap: 12px;
      align-content: center;
      justify-items: stretch;
      min-width: 0;
      z-index: 1;
    }}
    .media-column.media-count-2 {{
      grid-template-columns: repeat(2, minmax(0, 1fr));
      align-items: stretch;
    }}
    .media-column.media-count-many {{
      grid-template-columns: repeat(2, minmax(0, 1fr));
      align-items: stretch;
    }}
    .media-column.media-count-many .media-frame:first-child {{
      grid-column: 1 / -1;
    }}
    .media-frame {{
      width: 100%;
      min-height: 180px;
      display: grid;
      place-items: center;
      background: transparent;
      border: 0;
      border-radius: 0;
      overflow: visible;
    }}
    .media-count-2 .media-frame,
    .media-count-many .media-frame {{
      min-height: 180px;
    }}
    .media-frame img {{
      display: block;
      max-width: 100%;
      max-height: 520px;
      width: auto;
      height: auto;
      object-fit: contain;
    }}
    body.editing .media-frame img {{
      cursor: move;
      outline: 2px solid transparent;
      outline-offset: -2px;
    }}
    body.editing .media-frame img.selected-image {{
      outline-color: var(--accent);
      box-shadow: 0 0 0 4px rgba(11, 116, 222, 0.18);
    }}
    .media-frame.free-position {{
      position: absolute;
      width: 32%;
      min-height: 0;
      z-index: 4;
      display: block;
      touch-action: none;
    }}
    .media-frame.free-position img {{
      width: 100%;
      max-width: none;
      max-height: none;
    }}
    .media-frame .resize-handle {{
      display: none;
    }}
    body.editing .media-frame.free-position.selected-frame .resize-handle {{
      display: block;
    }}
    .resize-handle {{
      position: absolute;
      width: 12px;
      height: 12px;
      border: 2px solid #fff;
      border-radius: 999px;
      background: var(--accent);
      box-shadow: 0 1px 6px rgba(0,0,0,0.22);
      z-index: 8;
      touch-action: none;
    }}
    .resize-handle.nw {{ left: -6px; top: -6px; cursor: nwse-resize; }}
    .resize-handle.ne {{ right: -6px; top: -6px; cursor: nesw-resize; }}
    .resize-handle.sw {{ left: -6px; bottom: -6px; cursor: nesw-resize; }}
    .resize-handle.se {{ right: -6px; bottom: -6px; cursor: nwse-resize; }}
    .free-text-box {{
      position: absolute;
      left: 12%;
      top: 18%;
      width: 34%;
      min-height: 52px;
      padding: 10px 12px;
      color: var(--text);
      font: 700 24px/1.28 var(--font);
      background: rgba(255,255,255,0.7);
      border: 1px dashed var(--accent);
      border-radius: 7px;
      z-index: 5;
      cursor: move;
      touch-action: none;
      max-width: calc(100% - 32px);
      max-height: calc(100% - 32px);
      overflow: hidden;
      overflow-wrap: anywhere;
      word-break: break-word;
      contain: layout paint;
    }}
    body.editing .free-text-box {{
      outline: 2px solid transparent;
      outline-offset: 2px;
    }}
    body.editing .free-text-box.selected-text-box {{
      outline-color: var(--accent);
      box-shadow: 0 0 0 4px rgba(11, 116, 222, 0.14);
    }}
    .free-text-box .resize-handle {{
      display: none;
    }}
    body.editing .free-text-box.selected-text-box .resize-handle {{
      display: block;
    }}
    body.theme-doodle .free-text-box {{
      background: rgba(255, 241, 216, 0.72);
      border: 3px solid var(--border);
      border-radius: 4px 10px 5px 8px;
      box-shadow: 6px 7px 0 rgba(59, 41, 37, 0.12);
    }}
    .page-number {{
      position: absolute;
      right: 28px;
      bottom: 22px;
      font-size: 16px;
      color: var(--muted);
    }}
    body.editing .slide [data-editable="true"] {{
      outline: 2px dashed rgba(11, 116, 222, 0.35);
      outline-offset: 6px;
      cursor: text;
    }}
    @media (max-width: 900px) {{
      .deck-shell {{ padding: 12px; }}
      .deck-toolbar {{ flex-direction: column; align-items: stretch; }}
      .editor-toolbar {{ align-items: stretch; }}
      .editor-toolbar label {{ width: 100%; justify-content: space-between; }}
      .editor-toolbar input,
      .editor-toolbar select,
      .editor-toolbar button {{ max-width: 100%; }}
      .editor-toolbar .toolbar-divider {{ display: none; }}
      .deck-actions {{ justify-content: space-between; flex-wrap: wrap; }}
      .slide-stage {{ width: calc(100vw - 24px); }}
      .slide {{ aspect-ratio: auto; min-height: 78vh; }}
      .slide-inner {{
        grid-template-columns: 1fr;
        padding: 34px 26px 58px;
        gap: 24px;
      }}
      .slide-inner::before {{ left: 26px; bottom: 38px; }}
      .slide-title {{ font-size: 34px; }}
      .slide-body {{ font-size: 24px; }}
      .slide.text-only .slide-inner {{ padding-right: 26px; }}
      .media-frame {{ min-height: 220px; }}
      .media-column.media-count-2,
      .media-column.media-count-many {{ grid-template-columns: 1fr; }}
      .media-column.media-count-many .media-frame:first-child {{ grid-column: auto; }}
    }}
    @media print {{
      body {{ background: #fff; }}
      .deck-toolbar {{ display: none; }}
      .deck-shell {{ padding: 0; }}
      .slide-stage {{ width: 100%; }}
      .slide {{ display: grid !important; page-break-after: always; box-shadow: none; border: 0; }}
    }}
    body.embedded {{
      background: #fff;
      overflow: hidden;
    }}
    body.embedded .deck-shell {{
      min-height: 100vh;
      padding: 0;
      display: grid;
      place-items: center;
    }}
    body.embedded .deck-toolbar {{
      display: none;
    }}
    body.embedded .editor-toolbar {{
      display: none !important;
    }}
    body.embedded.editing .editor-toolbar {{
      display: flex !important;
      position: fixed;
      left: 10px;
      right: 10px;
      top: 8px;
      max-height: 44vh;
      overflow: auto;
      z-index: 100;
      box-shadow: 0 14px 30px rgba(15, 35, 70, 0.18);
    }}
    body.embedded.editing .slide-stage {{
      padding-top: 76px;
    }}
    body.embedded .slide-stage {{
      width: 100%;
      padding: 12px;
    }}
    body.embedded .slide {{
      border: 0;
      box-shadow: none;
    }}
    body.embedded .slide-inner {{
      padding: 38px 42px;
      gap: 20px;
    }}
    body.embedded .slide.text-only .slide-inner {{
      padding-right: 42px;
    }}
    body.embedded .media-frame {{
      min-height: 210px;
    }}
    body.embedded .page-number {{
      display: none;
    }}
    body.export-scroll {{
      overflow-x: hidden;
      overflow-y: auto;
    }}
    body.export-scroll .deck-shell {{
      min-height: 100vh;
      padding: 20px;
    }}
    body.export-scroll .deck-toolbar,
    body.export-scroll .editor-toolbar {{
      display: none !important;
    }}
    body.export-scroll .slide-stage {{
      width: min(1280px, calc(100vw - 40px));
      display: grid;
      gap: 20px;
      margin: 0 auto;
    }}
    body.export-scroll .slide {{
      display: grid !important;
      margin: 0 auto;
      page-break-after: always;
    }}
    body.export-scroll .slide:not(.active) {{
      display: grid !important;
    }}
  </style>
</head>
<body class="theme-{html.escape(style_key)}">
  <main class="deck-shell" data-deck="{data_blob}">
    <div class="deck-toolbar">
      <div class="deck-title">{html.escape(title)}</div>
      <div class="deck-actions">
        <button onclick="previousSlide()" aria-label="Previous slide">Previous</button>
        <span class="counter"><span id="currentSlide">1</span> / {slide_count}</span>
        <button onclick="nextSlide()" aria-label="Next slide">Next</button>
        <button onclick="toggleEdit()" id="editButton">Edit HTML</button>
        <button class="primary" onclick="downloadEdited('paged')">Download Paged HTML</button>
        <button onclick="downloadEdited('scroll')">Download Scroll HTML</button>
        <button onclick="saveEditedToServer(false)" id="saveEditedButton">Save Edits</button>
        <button onclick="saveEditedToServer(true)" id="downloadEditedZipButton">Download Edited ZIP</button>
      </div>
    </div>
    <div class="editor-toolbar" id="editorToolbar" aria-label="HTML editor toolbar">
      <button onclick="undoEdit()" id="undoButton" disabled>Undo</button>
      <button onclick="applyBold()" title="Bold selected text"><strong>Bold</strong></button>
      <button onclick="applyItalic()" title="Italic selected text"><em>Italic</em></button>
      <button onclick="applyUnderline()" title="Underline selected text"><u>Underline</u></button>
      <button onclick="applyTextAlign('left')" title="Align selected text left">Left</button>
      <button onclick="applyTextAlign('center')" title="Align selected text center">Center</button>
      <button onclick="applyTextAlign('right')" title="Align selected text right">Right</button>
      <label>Size <input id="editFontSize" type="number" min="12" max="96" value="28"></label>
      <button onclick="applyFontSize()">Apply Size</button>
      <label>Font
        <select id="editFontFamily">
          <option value="Inter, Segoe UI, Arial, sans-serif">Inter</option>
          <option value="Aptos, Segoe UI, Arial, sans-serif">Aptos</option>
          <option value="Calibri, Arial, sans-serif">Calibri</option>
          <option value="Arial, sans-serif">Arial</option>
          <option value="'Microsoft YaHei', 'PingFang SC', 'Noto Sans CJK SC', sans-serif">Microsoft YaHei</option>
          <option value="'SimHei', 'Microsoft YaHei', sans-serif">SimHei</option>
          <option value="'SimSun', 'Songti SC', serif">SimSun</option>
          <option value="'KaiTi', 'STKaiti', serif">KaiTi</option>
          <option value="'Noto Sans SC', 'Microsoft YaHei', sans-serif">Noto Sans SC</option>
          <option value="'Source Han Sans SC', 'Microsoft YaHei', sans-serif">Source Han Sans SC</option>
          <option value="'PingFang SC', 'Microsoft YaHei', sans-serif">PingFang SC</option>
          <option value="'HarmonyOS Sans SC', 'Microsoft YaHei', sans-serif">HarmonyOS Sans</option>
          <option value="Georgia, serif">Georgia</option>
          <option value="'Times New Roman', serif">Times New Roman</option>
          <option value="'Trebuchet MS', Arial, sans-serif">Trebuchet MS</option>
          <option value="Verdana, Geneva, sans-serif">Verdana</option>
          <option value="Tahoma, Geneva, sans-serif">Tahoma</option>
          <option value="'Segoe Print', 'Comic Sans MS', cursive">Segoe Print</option>
          <option value="'Comic Sans MS', 'Segoe Print', cursive">Comic Sans MS</option>
          <option value="'Ink Free', 'Segoe Print', cursive">Ink Free</option>
          <option value="'Courier New', monospace">Courier New</option>
          <option value="'Cascadia Mono', 'Courier New', monospace">Cascadia Mono</option>
        </select>
      </label>
      <button onclick="applyFontFamily()">Apply Font</button>
      <span class="color-palette" aria-label="Preset text colors">
        <button class="color-swatch" style="--swatch:#000000" title="Black" onclick="setTextColor('#000000')"></button>
        <button class="color-swatch white" style="--swatch:#ffffff" title="White" onclick="setTextColor('#ffffff')"></button>
        <button class="color-swatch" style="--swatch:#ef4444" title="Red" onclick="setTextColor('#ef4444')"></button>
        <button class="color-swatch" style="--swatch:#f97316" title="Orange" onclick="setTextColor('#f97316')"></button>
        <button class="color-swatch" style="--swatch:#facc15" title="Yellow" onclick="setTextColor('#facc15')"></button>
        <button class="color-swatch" style="--swatch:#22c55e" title="Green" onclick="setTextColor('#22c55e')"></button>
        <button class="color-swatch" style="--swatch:#2563eb" title="Blue" onclick="setTextColor('#2563eb')"></button>
        <button class="color-swatch" style="--swatch:#7c3aed" title="Purple" onclick="setTextColor('#7c3aed')"></button>
      </span>
      <label>Custom <input id="editColor" type="color" value="{preset["accent"]}" oninput="setTextColor(this.value, false)"></label>
      <button onclick="applyTextColor()">Apply Color</button>
      <span class="toolbar-divider"></span>
      <label>Image width % <input id="imageWidth" type="number" min="10" max="100" value="100"></label>
      <button onclick="applyImageWidth()">Resize Image</button>
      <button onclick="deleteSelectedImage()">Delete Image</button>
      <label>Add image URL <input id="newImageUrl" type="url" placeholder="https://... or assets/name.png"></label>
      <button onclick="addImageFromUrl()">Add URL Image</button>
      <label><input id="addImageFile" type="file" accept="image/*" onchange="addImageFromFile(this)">Add Local Image</label>
      <button onclick="addTextBox()">Add Text Box</button>
      <button onclick="addTable()">Add Table</button>
      <button onclick="convertSelectionToTable()">Auto Table</button>
      <span class="edit-hint">Drag images or added text boxes; drag text-box borders/corners to resize</span>
    </div>
    <section class="slide-stage">
      {slide_markup}
    </section>
  </main>
  <script>
    if (window.self !== window.top) {{
      document.body.classList.add('embedded');
    }}
    const slides = Array.from(document.querySelectorAll('.slide'));
    const deckJobId = {json.dumps(job_id)};
    let current = 0;
    let savedRange = null;
    let selectedImage = null;
    let selectedTextBox = null;
    let dragState = null;
    const editHistory = [];
    let pendingHistory = null;
    document.addEventListener('selectionchange', () => {{
      const selection = window.getSelection();
      if (!document.body.classList.contains('editing') || !selection || selection.rangeCount === 0) return;
      const node = selection.anchorNode;
      if (node && document.querySelector('.slide.active')?.contains(node)) {{
        savedRange = selection.getRangeAt(0).cloneRange();
      }}
    }});
    document.getElementById('editorToolbar').addEventListener('mousedown', event => {{
      if (event.target.tagName !== 'INPUT' && event.target.tagName !== 'SELECT') {{
        event.preventDefault();
      }}
    }});
    function showSlide(index) {{
      current = Math.max(0, Math.min(slides.length - 1, index));
      slides.forEach((slide, i) => slide.classList.toggle('active', i === current));
      document.getElementById('currentSlide').textContent = String(current + 1);
      requestAnimationFrame(() => fitSlideText(slides[current]));
      clearSelectedImage();
      clearSelectedTextBox();
      updateUndoButton();
      location.hash = 'slide-' + (current + 1);
    }}
    function textOverflowing(slide) {{
      const inner = slide.querySelector('.slide-inner');
      const textColumn = slide.querySelector('.text-column');
      const body = slide.querySelector('.slide-body');
      const title = slide.querySelector('.slide-title');
      const targets = [inner, textColumn, body, title, ...slide.querySelectorAll('.body-card, .content-list, .timeline-step, .free-text-box')].filter(Boolean);
      return targets.some(el => {{
        const heightOverflow = el.scrollHeight > el.clientHeight + 2 && el.clientHeight > 0;
        const widthOverflow = el.scrollWidth > el.clientWidth + 2 && el.clientWidth > 0;
        return heightOverflow || widthOverflow;
      }});
    }}
    function shrinkTextTargets(slide) {{
      slide.querySelectorAll('.slide-title, .slide-body, .body-card, .content-list, .timeline-step, .free-text-box').forEach(el => {{
        const computed = window.getComputedStyle(el);
        const size = parseFloat(computed.fontSize) || 24;
        const lineHeight = parseFloat(computed.lineHeight) || size * 1.3;
        const isTitle = el.classList.contains('slide-title');
        const minSize = isTitle ? 24 : 20;
        const nextSize = Math.max(minSize, size * 0.93);
        el.style.fontSize = nextSize.toFixed(2) + 'px';
        el.style.lineHeight = Math.max(nextSize * 1.08, lineHeight * 0.96).toFixed(2) + 'px';
      }});
    }}
    function resetAutoFit(slide) {{
      slide.querySelectorAll('[data-auto-fit="true"]').forEach(el => {{
        el.style.fontSize = '';
        el.style.lineHeight = '';
        delete el.dataset.autoFit;
      }});
    }}
    function markAutoFit(slide) {{
      slide.querySelectorAll('.slide-title, .slide-body, .body-card, .content-list, .timeline-step, .free-text-box').forEach(el => {{
        if (el.style.fontSize || el.style.lineHeight) el.dataset.autoFit = 'true';
      }});
    }}
    function fitSlideText(slide) {{
      if (!slide) return;
      resetAutoFit(slide);
      for (let i = 0; i < 16 && textOverflowing(slide); i += 1) {{
        shrinkTextTargets(slide);
      }}
      markAutoFit(slide);
    }}
    function fitAllSlides() {{
      const original = current;
      slides.forEach((slide, index) => {{
        slides.forEach((item, i) => item.classList.toggle('active', i === index));
        fitSlideText(slide);
      }});
      slides.forEach((item, i) => item.classList.toggle('active', i === original));
      document.getElementById('currentSlide').textContent = String(original + 1);
    }}
    window.addEventListener('resize', () => fitSlideText(slides[current]));
    function nextSlide() {{ showSlide(current + 1); }}
    function previousSlide() {{ showSlide(current - 1); }}
    function toggleEdit() {{
      const editing = !document.body.classList.contains('editing');
      document.body.classList.toggle('editing', editing);
      document.querySelectorAll('[data-editable="true"]').forEach(el => {{
        el.contentEditable = editing ? 'true' : 'false';
      }});
      if (!editing) {{
        clearSelectedImage();
        clearSelectedTextBox();
      }}
      bindEditableImages();
      bindDraggableTextBoxes();
      document.getElementById('editButton').textContent = editing ? 'Stop Editing' : 'Edit HTML';
    }}
    function restoreSelection() {{
      if (!savedRange) return false;
      const selection = window.getSelection();
      selection.removeAllRanges();
      selection.addRange(savedRange);
      return true;
    }}
    function activeSlideInner() {{
      return document.querySelector('.slide.active .slide-inner');
    }}
    function beginHistory(label) {{
      const slideInner = activeSlideInner();
      if (!slideInner) return null;
      pendingHistory = {{
        label,
        slideIndex: current,
        before: slideInner.innerHTML,
      }};
      return pendingHistory;
    }}
    function commitHistory(entry = pendingHistory) {{
      if (!entry) return;
      const slide = slides[entry.slideIndex];
      const slideInner = slide?.querySelector('.slide-inner');
      const after = slideInner?.innerHTML || '';
      if (after && after !== entry.before) {{
        editHistory.push({{ ...entry, after }});
        if (editHistory.length > 60) editHistory.shift();
      }}
      pendingHistory = null;
      updateUndoButton();
    }}
    function cancelHistory(entry = pendingHistory) {{
      if (entry === pendingHistory) pendingHistory = null;
      updateUndoButton();
    }}
    function updateUndoButton() {{
      const button = document.getElementById('undoButton');
      if (button) button.disabled = editHistory.length === 0;
    }}
    function undoEdit() {{
      const entry = editHistory.pop();
      if (!entry) return;
      const slide = slides[entry.slideIndex];
      const slideInner = slide?.querySelector('.slide-inner');
      if (!slideInner) return;
      current = entry.slideIndex;
      slides.forEach((item, i) => item.classList.toggle('active', i === current));
      document.getElementById('currentSlide').textContent = String(current + 1);
      slideInner.innerHTML = entry.before;
      slideInner.querySelectorAll('[data-bound-image-editor], [data-bound-text-drag]').forEach(el => {{
        delete el.dataset.boundImageEditor;
        delete el.dataset.boundTextDrag;
      }});
      clearSelectedImage();
      clearSelectedTextBox();
      bindEditableImages();
      bindDraggableTextBoxes();
      updateUndoButton();
      location.hash = 'slide-' + (current + 1);
    }}
    function wrapSelection(styleMap) {{
      if (!restoreSelection()) return;
      const selection = window.getSelection();
      if (!selection || selection.rangeCount === 0 || selection.isCollapsed) return;
      const historyEntry = beginHistory('text style');
      const range = selection.getRangeAt(0);
      const span = document.createElement('span');
      Object.assign(span.style, styleMap);
      span.appendChild(range.extractContents());
      range.insertNode(span);
      range.selectNodeContents(span);
      selection.removeAllRanges();
      selection.addRange(range);
      savedRange = range.cloneRange();
      commitHistory(historyEntry);
    }}
    function applyFontSize() {{
      const px = Math.max(12, Math.min(96, parseInt(document.getElementById('editFontSize').value, 10) || 28));
      wrapSelection({{ fontSize: px + 'px' }});
    }}
    function applyBold() {{
      wrapSelection({{ fontWeight: '800' }});
    }}
    function applyItalic() {{
      wrapSelection({{ fontStyle: 'italic' }});
    }}
    function applyUnderline() {{
      wrapSelection({{ textDecoration: 'underline' }});
    }}
    function applyTextAlign(alignment) {{
      const allowed = ['left', 'center', 'right'];
      if (!allowed.includes(alignment)) return;
      const historyEntry = beginHistory('text align');
      let changed = false;
      if (selectedTextBox) {{
        selectedTextBox.style.textAlign = alignment;
        changed = true;
      }}
      if (restoreSelection()) {{
        const selection = window.getSelection();
        if (selection && selection.rangeCount > 0) {{
          const range = selection.getRangeAt(0);
          const container = range.commonAncestorContainer.nodeType === 1 ? range.commonAncestorContainer : range.commonAncestorContainer.parentElement;
          const root = document.querySelector('.slide.active');
          const candidates = Array.from(root?.querySelectorAll('p, li, h1, .body-card, .free-text-box, .slide-body') || []);
          let targets = candidates.filter(el => range.intersectsNode(el));
          if (!targets.length && container) {{
            const fallback = container.closest?.('p, li, h1, .body-card, .free-text-box, .slide-body, [data-editable="true"]');
            if (fallback) targets = [fallback];
          }}
          targets.forEach(el => {{
            el.style.textAlign = alignment;
            changed = true;
          }});
        }}
      }}
      if (changed) commitHistory(historyEntry); else cancelHistory(historyEntry);
    }}
    function applyFontFamily() {{
      wrapSelection({{ fontFamily: document.getElementById('editFontFamily').value }});
    }}
    function setTextColor(color, applyNow = true) {{
      document.getElementById('editColor').value = color;
      if (applyNow) wrapSelection({{ color }});
    }}
    function applyTextColor() {{
      setTextColor(document.getElementById('editColor').value);
    }}
    function bindEditableImages() {{
      document.querySelectorAll('.media-frame img').forEach(img => {{
        if (img.dataset.boundImageEditor === 'true') return;
        img.dataset.boundImageEditor = 'true';
        img.addEventListener('pointerdown', event => {{
          if (!document.body.classList.contains('editing')) return;
          event.preventDefault();
          event.stopPropagation();
          selectImage(img);
          startImageDrag(img, event);
        }});
      }});
    }}
    function bindDraggableTextBoxes() {{
      document.querySelectorAll('.free-text-box').forEach(box => {{
        if (box.dataset.boundTextDrag === 'true') return;
        box.dataset.boundTextDrag = 'true';
        box.addEventListener('pointerdown', event => {{
          if (!document.body.classList.contains('editing')) return;
          if (event.target.classList?.contains('resize-handle')) return;
          selectTextBox(box);
          if (event.target !== box) return;
          if (isTextBoxBorderPointer(box, event)) {{
            event.preventDefault();
            event.stopPropagation();
            startTextResizeDrag(box, inferTextResizeCorner(box, event), event);
            return;
          }}
          startElementDrag(box, event);
        }});
      }});
    }}
    function selectTextBox(box) {{
      clearSelectedImage();
      if (selectedTextBox && selectedTextBox !== box) selectedTextBox.classList.remove('selected-text-box');
      selectedTextBox = box;
      box.classList.add('selected-text-box');
      ensureTextResizeHandles(box);
    }}
    function clearSelectedTextBox() {{
      document.querySelectorAll('.selected-text-box').forEach(box => box.classList.remove('selected-text-box'));
      selectedTextBox = null;
    }}
    function ensureTextResizeHandles(box) {{
      if (!box || box.querySelector('.resize-handle')) return;
      ['nw', 'ne', 'sw', 'se'].forEach(corner => {{
        const handle = document.createElement('span');
        handle.className = 'resize-handle text-resize-handle ' + corner;
        handle.dataset.corner = corner;
        handle.title = 'Drag to resize text box';
        handle.addEventListener('pointerdown', event => {{
          if (!document.body.classList.contains('editing')) return;
          event.preventDefault();
          event.stopPropagation();
          selectTextBox(box);
          startTextResizeDrag(box, corner, event);
        }});
        box.appendChild(handle);
      }});
    }}
    function isTextBoxBorderPointer(box, event) {{
      const rect = box.getBoundingClientRect();
      const edge = 12;
      return (
        event.clientX - rect.left <= edge ||
        rect.right - event.clientX <= edge ||
        event.clientY - rect.top <= edge ||
        rect.bottom - event.clientY <= edge
      );
    }}
    function inferTextResizeCorner(box, event) {{
      const rect = box.getBoundingClientRect();
      const horizontal = event.clientX < rect.left + rect.width / 2 ? 'w' : 'e';
      const vertical = event.clientY < rect.top + rect.height / 2 ? 'n' : 's';
      return vertical + horizontal;
    }}
    function selectImage(img) {{
      clearSelectedImage();
      clearSelectedTextBox();
      selectedImage = img;
      img.classList.add('selected-image');
      const frame = img.closest('.media-frame');
      frame?.classList.add('selected-frame');
      ensureResizeHandles(frame);
      const widthValue = img.closest('.free-position') ? parseInt(img.closest('.free-position').style.width, 10) : parseInt(img.style.width, 10);
      document.getElementById('imageWidth').value = widthValue || 100;
    }}
    function ensureResizeHandles(frame) {{
      if (!frame || frame.querySelector('.resize-handle')) return;
      ['nw', 'ne', 'sw', 'se'].forEach(corner => {{
        const handle = document.createElement('span');
        handle.className = 'resize-handle ' + corner;
        handle.dataset.corner = corner;
        handle.addEventListener('pointerdown', event => {{
          if (!document.body.classList.contains('editing')) return;
          event.preventDefault();
          event.stopPropagation();
          startResizeDrag(frame, corner, event);
        }});
        frame.appendChild(handle);
      }});
    }}
    function clearSelectedImage() {{
      document.querySelectorAll('.selected-image').forEach(img => img.classList.remove('selected-image'));
      document.querySelectorAll('.selected-frame').forEach(frame => frame.classList.remove('selected-frame'));
      selectedImage = null;
    }}
    function slideInnerFor(element) {{
      return element.closest('.slide')?.querySelector('.slide-inner');
    }}
    function makeFrameFree(frame, img) {{
      if (frame.classList.contains('free-position')) return;
      const slideInner = slideInnerFor(frame);
      const oldColumn = frame.closest('.media-column');
      if (!slideInner) return;
      const imageRect = img.getBoundingClientRect();
      const innerRect = slideInner.getBoundingClientRect();
      const widthPct = Math.max(6, Math.min(78, imageRect.width / innerRect.width * 100));
      const leftPct = Math.max(0, Math.min(100 - widthPct, (imageRect.left - innerRect.left) / innerRect.width * 100));
      const topPct = Math.max(0, Math.min(92, (imageRect.top - innerRect.top) / innerRect.height * 100));
      if (oldColumn) oldColumn.dataset.layoutAnchor = 'true';
      slideInner.insertBefore(frame, slideInner.querySelector('.page-number'));
      frame.classList.add('free-position');
      frame.style.left = leftPct + '%';
      frame.style.top = topPct + '%';
      frame.style.width = widthPct + '%';
      img.style.width = '100%';
      ensureResizeHandles(frame);
    }}
    function startImageDrag(img, event) {{
      const frame = img.closest('.media-frame');
      const historyEntry = beginHistory('image move');
      makeFrameFree(frame, img);
      startElementDrag(frame, event, historyEntry);
    }}
    function startElementDrag(element, event, historyEntry = null) {{
      const slideInner = slideInnerFor(element);
      if (!slideInner) return;
      const entry = historyEntry || beginHistory(element.classList.contains('free-text-box') ? 'text box move' : 'element move');
      const elementRect = element.getBoundingClientRect();
      const innerRect = slideInner.getBoundingClientRect();
      dragState = {{
        mode: 'move',
        element,
        historyEntry: entry,
        slideInner,
        startX: event.clientX,
        startY: event.clientY,
        startLeft: elementRect.left - innerRect.left,
        startTop: elementRect.top - innerRect.top,
        width: elementRect.width,
        height: elementRect.height,
      }};
      element.setPointerCapture?.(event.pointerId);
    }}
    function startResizeDrag(frame, corner, event) {{
      const slideInner = slideInnerFor(frame);
      if (!slideInner) return;
      const historyEntry = beginHistory('image resize');
      const frameRect = frame.getBoundingClientRect();
      const innerRect = slideInner.getBoundingClientRect();
      dragState = {{
        mode: 'resize',
        element: frame,
        historyEntry,
        slideInner,
        corner,
        startX: event.clientX,
        startY: event.clientY,
        startLeft: frameRect.left - innerRect.left,
        startTop: frameRect.top - innerRect.top,
        width: frameRect.width,
        height: frameRect.height,
        aspect: frameRect.width / Math.max(1, frameRect.height),
      }};
      frame.setPointerCapture?.(event.pointerId);
    }}
    function startTextResizeDrag(box, corner, event) {{
      const slideInner = slideInnerFor(box);
      if (!slideInner) return;
      const historyEntry = beginHistory('text box resize');
      const boxRect = box.getBoundingClientRect();
      const innerRect = slideInner.getBoundingClientRect();
      dragState = {{
        mode: 'text-resize',
        element: box,
        historyEntry,
        slideInner,
        corner,
        startX: event.clientX,
        startY: event.clientY,
        startLeft: boxRect.left - innerRect.left,
        startTop: boxRect.top - innerRect.top,
        width: boxRect.width,
        height: boxRect.height,
      }};
      box.setPointerCapture?.(event.pointerId);
    }}
    function handleDragMove(event) {{
      if (!dragState) return;
      event.preventDefault();
      const innerRect = dragState.slideInner.getBoundingClientRect();
      if (dragState.mode === 'resize') {{
        const dx = event.clientX - dragState.startX;
        const east = dragState.corner.includes('e');
        const south = dragState.corner.includes('s');
        let width = east ? dragState.width + dx : dragState.width - dx;
        width = Math.max(40, Math.min(innerRect.width * 0.86, width));
        const height = width / Math.max(0.1, dragState.aspect);
        let left = east ? dragState.startLeft : dragState.startLeft + (dragState.width - width);
        let top = south ? dragState.startTop : dragState.startTop + (dragState.height - height);
        left = Math.max(0, Math.min(innerRect.width - width, left));
        top = Math.max(0, Math.min(innerRect.height - height, top));
        dragState.element.style.left = (left / innerRect.width * 100) + '%';
        dragState.element.style.top = (top / innerRect.height * 100) + '%';
        dragState.element.style.width = (width / innerRect.width * 100) + '%';
        document.getElementById('imageWidth').value = Math.round(width / innerRect.width * 100);
        return;
      }}
      if (dragState.mode === 'text-resize') {{
        const dx = event.clientX - dragState.startX;
        const dy = event.clientY - dragState.startY;
        const east = dragState.corner.includes('e');
        const south = dragState.corner.includes('s');
        const minWidth = Math.min(180, innerRect.width * 0.35);
        const minHeight = 56;
        let width = east ? dragState.width + dx : dragState.width - dx;
        let height = south ? dragState.height + dy : dragState.height - dy;
        width = Math.max(minWidth, Math.min(innerRect.width * 0.92, width));
        height = Math.max(minHeight, Math.min(innerRect.height * 0.82, height));
        let left = east ? dragState.startLeft : dragState.startLeft + (dragState.width - width);
        let top = south ? dragState.startTop : dragState.startTop + (dragState.height - height);
        left = Math.max(0, Math.min(innerRect.width - width, left));
        top = Math.max(0, Math.min(innerRect.height - height, top));
        dragState.element.style.left = (left / innerRect.width * 100) + '%';
        dragState.element.style.top = (top / innerRect.height * 100) + '%';
        dragState.element.style.width = (width / innerRect.width * 100) + '%';
        dragState.element.style.height = (height / innerRect.height * 100) + '%';
        dragState.element.style.minHeight = '0';
        return;
      }}
      const dragSensitivity = 0.82;
      const maxLeft = Math.max(0, innerRect.width - dragState.width);
      const maxTop = Math.max(0, innerRect.height - dragState.height);
      const left = Math.max(0, Math.min(maxLeft, dragState.startLeft + (event.clientX - dragState.startX) * dragSensitivity));
      const top = Math.max(0, Math.min(maxTop, dragState.startTop + (event.clientY - dragState.startY) * dragSensitivity));
      dragState.element.style.left = (left / innerRect.width * 100) + '%';
      dragState.element.style.top = (top / innerRect.height * 100) + '%';
    }}
    function finishDrag() {{
      if (dragState?.historyEntry) {{
        commitHistory(dragState.historyEntry);
      }}
      dragState = null;
    }}
    document.addEventListener('pointermove', handleDragMove);
    document.addEventListener('pointerup', finishDrag);
    function currentMediaColumn() {{
      const active = document.querySelector('.slide.active');
      let column = active.querySelector('.media-column');
      if (!column) {{
        column = document.createElement('div');
        column.className = 'media-column media-count-1';
        active.querySelector('.slide-inner').insertBefore(column, active.querySelector('.page-number'));
        active.classList.remove('text-only');
        active.classList.add('has-media');
      }}
      return column;
    }}
    function refreshMediaCount(column) {{
      const count = column.querySelectorAll('.media-frame').length;
      if (count === 0 && column.dataset.layoutAnchor === 'true') return;
      column.classList.remove('media-count-1', 'media-count-2', 'media-count-many');
      column.classList.add(count > 2 ? 'media-count-many' : count === 2 ? 'media-count-2' : 'media-count-1');
      const slide = column.closest('.slide');
      if (count === 0) {{
        column.remove();
        slide.classList.remove('has-media');
        slide.classList.add('text-only');
      }}
    }}
    function addImage(src) {{
      if (!src) return;
      const historyEntry = beginHistory('add image');
      const column = currentMediaColumn();
      delete column.dataset.layoutAnchor;
      const frame = document.createElement('figure');
      frame.className = 'media-frame';
      const img = document.createElement('img');
      img.src = src;
      img.alt = 'Added image';
      frame.appendChild(img);
      column.appendChild(frame);
      refreshMediaCount(column);
      bindEditableImages();
      clearSelectedImage();
      selectImage(img);
      commitHistory(historyEntry);
    }}
    function applyImageWidth() {{
      if (!selectedImage) return;
      const historyEntry = beginHistory('image width');
      const width = Math.max(10, Math.min(100, parseInt(document.getElementById('imageWidth').value, 10) || 100));
      const frame = selectedImage.closest('.media-frame');
      if (frame?.classList.contains('free-position')) {{
        frame.style.width = width + '%';
        selectedImage.style.width = '100%';
      }} else {{
        selectedImage.style.width = width + '%';
      }}
      commitHistory(historyEntry);
    }}
    function deleteSelectedImage() {{
      if (!selectedImage) return;
      const historyEntry = beginHistory('delete image');
      const column = selectedImage.closest('.media-column');
      selectedImage.closest('.media-frame').remove();
      selectedImage = null;
      if (column) refreshMediaCount(column);
      commitHistory(historyEntry);
    }}
    function addImageFromUrl() {{
      const input = document.getElementById('newImageUrl');
      addImage(input.value.trim());
      input.value = '';
    }}
    function addImageFromFile(input) {{
      const file = input.files && input.files[0];
      if (!file) return;
      const reader = new FileReader();
      reader.onload = () => addImage(reader.result);
      reader.readAsDataURL(file);
      input.value = '';
    }}
    function addTextBox() {{
      const active = document.querySelector('.slide.active');
      const slideInner = active?.querySelector('.slide-inner');
      if (!slideInner) return;
      const historyEntry = beginHistory('add text box');
      const box = document.createElement('div');
      box.className = 'free-text-box';
      box.dataset.editable = 'true';
      box.contentEditable = document.body.classList.contains('editing') ? 'true' : 'false';
      box.textContent = 'New text';
      slideInner.insertBefore(box, slideInner.querySelector('.page-number'));
      bindDraggableTextBoxes();
      selectTextBox(box);
      box.focus();
      const range = document.createRange();
      range.selectNodeContents(box);
      const selection = window.getSelection();
      selection.removeAllRanges();
      selection.addRange(range);
      savedRange = range.cloneRange();
      commitHistory(historyEntry);
    }}
    function parseTableRows(text) {{
      const lines = String(text || '').split(/\\r?\\n/).map(line => line.trim()).filter(Boolean);
      if (lines.length < 2) return null;
      let delimiter = '';
      if (lines.every(line => line.includes('|'))) delimiter = '|';
      else if (lines.every(line => line.includes('\\t'))) delimiter = '\\t';
      else if (lines.every(line => line.includes(','))) delimiter = ',';
      else return null;
      const rows = lines.map(line => line.split(delimiter).map(cell => cell.trim()));
      if (rows.some(row => row.length < 2)) return null;
      const maxCols = Math.max(...rows.map(row => row.length));
      const minCols = Math.min(...rows.map(row => row.length));
      if (maxCols - minCols > 1) return null;
      return rows.map(row => row.concat(Array(Math.max(0, maxCols - row.length)).fill('')));
    }}
    function buildTable(rows) {{
      const wrapper = document.createElement('div');
      wrapper.className = 'table-scroll';
      const table = document.createElement('table');
      table.className = 'data-table';
      const thead = document.createElement('thead');
      const headerRow = document.createElement('tr');
      rows[0].forEach(cell => {{
        const th = document.createElement('th');
        th.textContent = cell;
        headerRow.appendChild(th);
      }});
      thead.appendChild(headerRow);
      const tbody = document.createElement('tbody');
      rows.slice(1).forEach(row => {{
        const tr = document.createElement('tr');
        row.forEach(cell => {{
          const td = document.createElement('td');
          td.textContent = cell;
          tr.appendChild(td);
        }});
        tbody.appendChild(tr);
      }});
      table.appendChild(thead);
      table.appendChild(tbody);
      wrapper.appendChild(table);
      return wrapper;
    }}
    function addTable() {{
      const body = document.querySelector('.slide.active .slide-body');
      if (!body) return;
      const historyEntry = beginHistory('add table');
      const table = buildTable([
        ['Column A', 'Column B', 'Column C'],
        ['Value 1', 'Value 2', 'Value 3'],
        ['Value 4', 'Value 5', 'Value 6'],
      ]);
      body.appendChild(table);
      commitHistory(historyEntry);
    }}
    function convertSelectionToTable() {{
      const historyEntry = beginHistory('auto table');
      let changed = false;
      if (restoreSelection()) {{
        const selection = window.getSelection();
        const selectedText = selection?.toString() || '';
        const rows = parseTableRows(selectedText);
        if (rows && selection.rangeCount) {{
          const range = selection.getRangeAt(0);
          range.deleteContents();
          range.insertNode(buildTable(rows));
          selection.removeAllRanges();
          savedRange = null;
          changed = true;
        }}
      }}
      if (!changed) {{
        const paragraphs = Array.from(document.querySelectorAll('.slide.active .slide-body p'));
        paragraphs.forEach(paragraph => {{
          const rows = parseTableRows(paragraph.innerText);
          if (!rows) return;
          paragraph.replaceWith(buildTable(rows));
          changed = true;
        }});
      }}
      if (changed) commitHistory(historyEntry);
      else cancelHistory(historyEntry);
    }}
    async function inlineLocalImages(root) {{
      const images = Array.from(root.querySelectorAll('img'));
      await Promise.all(images.map(async img => {{
        const src = img.getAttribute('src') || '';
        if (!src || /^(data:|https?:|blob:)/i.test(src)) return;
        try {{
          const response = await fetch(src);
          if (!response.ok) return;
          const blob = await response.blob();
          const dataUrl = await new Promise(resolve => {{
            const reader = new FileReader();
            reader.onload = () => resolve(reader.result);
            reader.readAsDataURL(blob);
          }});
          img.setAttribute('src', dataUrl);
        }} catch (error) {{
          console.warn('Could not inline image for offline download', src, error);
        }}
      }}));
    }}
    function cleanExportClone(clone, mode) {{
      clone.querySelector('#editorToolbar')?.remove();
      clone.querySelectorAll('[contenteditable]').forEach(el => el.removeAttribute('contenteditable'));
      clone.querySelectorAll('.resize-handle').forEach(el => el.remove());
      clone.querySelectorAll('.selected-image').forEach(el => el.classList.remove('selected-image'));
      clone.querySelectorAll('.selected-frame').forEach(el => el.classList.remove('selected-frame'));
      clone.querySelectorAll('.selected-text-box').forEach(el => el.classList.remove('selected-text-box'));
      const body = clone.querySelector('body');
      body?.classList.remove('editing', 'embedded');
      if (mode === 'scroll') {{
        body?.classList.add('export-scroll');
        clone.querySelector('.deck-toolbar')?.remove();
        clone.querySelectorAll('.slide').forEach(slide => slide.classList.add('active'));
      }} else {{
        body?.classList.remove('export-scroll');
      }}
    }}
    async function downloadEdited(mode = 'paged') {{
      const html = await exportEditedHtml(mode);
      const blob = new Blob([html], {{ type: 'text/html;charset=utf-8' }});
      const link = document.createElement('a');
      link.href = URL.createObjectURL(blob);
      link.download = mode === 'scroll' ? 'optimized-ppt-scroll.html' : 'optimized-ppt-paged.html';
      link.click();
      URL.revokeObjectURL(link.href);
    }}
    async function exportEditedHtml(mode = 'paged') {{
      if (document.activeElement && document.activeElement.blur) document.activeElement.blur();
      fitAllSlides();
      const clone = document.documentElement.cloneNode(true);
      cleanExportClone(clone, mode);
      await inlineLocalImages(clone);
      return '<!DOCTYPE html>\\n' + clone.outerHTML;
    }}
    async function saveEditedToServer(downloadZip = false) {{
      if (!deckJobId) {{
        alert('This file was generated before editable ZIP saving was added. Please regenerate it from the platform.');
        return false;
      }}
      const button = downloadZip ? document.getElementById('downloadEditedZipButton') : document.getElementById('saveEditedButton');
      const oldText = button?.textContent;
      if (button) {{
        button.disabled = true;
        button.textContent = 'Saving edits...';
      }}
      try {{
        const response = await fetch('/api/jobs/' + encodeURIComponent(deckJobId) + '/save-edited', {{
          method: 'POST',
          headers: {{ 'Content-Type': 'application/json' }},
          body: JSON.stringify({{
            pagedHtml: await exportEditedHtml('paged'),
            scrollHtml: await exportEditedHtml('scroll'),
          }}),
        }});
        const data = await response.json().catch(() => ({{}}));
        if (!response.ok) throw new Error(data.message || data.error || 'Could not save edited HTML');
        if (downloadZip) window.location.href = '/api/jobs/' + encodeURIComponent(deckJobId) + '/download?v=' + Date.now();
        else alert('Edited HTML has been saved. You can now download the ZIP package from the platform.');
        return true;
      }} catch (error) {{
        alert(error.message || 'Could not save edited HTML');
        return false;
      }} finally {{
        if (button) {{
          button.disabled = false;
          button.textContent = oldText || (downloadZip ? 'Download Edited ZIP' : 'Save Edits');
        }}
      }}
    }}
    window.exportEditedHtml = exportEditedHtml;
    window.saveEditedToServer = saveEditedToServer;
    document.addEventListener('keydown', event => {{
      const target = event.target;
      const isTyping = target?.closest?.('[contenteditable="true"], input, textarea, select');
      if ((event.ctrlKey || event.metaKey) && event.key.toLowerCase() === 'z' && document.body.classList.contains('editing') && !isTyping) {{
        event.preventDefault();
        undoEdit();
        return;
      }}
      if (event.key === 'ArrowRight') nextSlide();
      if (event.key === 'ArrowLeft') previousSlide();
    }});
    const hashMatch = location.hash.match(/slide-(\\d+)/);
    showSlide(hashMatch ? Number(hashMatch[1]) - 1 : 0);
  </script>
</body>
</html>"""


def parse_table_rows_from_text(content: str) -> List[List[str]]:
    lines = [line.strip() for line in str(content or "").splitlines() if line.strip()]
    if len(lines) < 2:
        return []
    delimiter = ""
    if all("|" in line for line in lines):
        delimiter = "|"
    elif all("\t" in line for line in lines):
        delimiter = "\t"
    elif all("," in line for line in lines):
        delimiter = ","
    if not delimiter:
        return []

    rows: List[List[str]] = []
    for line in lines:
        cells = [cell.strip() for cell in line.split(delimiter)]
        if len(cells) < 2:
            return []
        rows.append(cells)
    column_counts = {len(row) for row in rows}
    if len(column_counts) > 2:
        return []
    max_cols = max(column_counts)
    return [row + [""] * (max_cols - len(row)) for row in rows]


def render_table(rows: List[List[str]]) -> str:
    if not rows:
        return ""
    header = "".join(f"<th>{html.escape(cell)}</th>" for cell in rows[0])
    body_rows = rows[1:] if len(rows) > 1 else []
    body = "\n".join(
        "<tr>" + "".join(f"<td>{html.escape(cell)}</td>" for cell in row) + "</tr>"
        for row in body_rows
    )
    return f"""<div class="table-scroll">
          <table class="data-table">
            <thead><tr>{header}</tr></thead>
            <tbody>{body}</tbody>
          </table>
        </div>"""


def render_inline_text(content: str) -> str:
    chapter_match = re.match(r"^(chapter\s+[0-9ivxlcdm]+)(.*)$", content.strip(), flags=re.IGNORECASE)
    if chapter_match:
        label = html.escape(chapter_match.group(1).upper())
        rest = html.escape(chapter_match.group(2).strip())
        if rest:
            return f'<span class="chapter-label">{label}</span><span class="chapter-rest"> {rest}</span>'
        return f'<span class="chapter-label">{label}</span>'
    escaped = html.escape(content)
    match = re.match(r"^([^:\n]{1,42}:)(\s+.+)$", content)
    if not match or "://" in match.group(1):
        return escaped
    prefix = html.escape(match.group(1))
    rest = html.escape(match.group(2))
    return f'<span class="accent">{prefix}</span>{rest}'


def should_render_as_list(lines: List[str]) -> bool:
    if len(lines) < 2 or len(lines) > 10:
        return False
    average_length = sum(len(line) for line in lines) / max(1, len(lines))
    if average_length > 120:
        return False
    return True


def render_body_item(item: Dict[str, Any], index: int = 0) -> str:
    content = str(item.get("content") or "").strip()
    if not content:
        return ""
    rows = normalize_table_rows(item.get("rows"))
    if not rows and item.get("kind") == "table":
        rows = parse_table_rows_from_text(content)
    if not rows:
        rows = parse_table_rows_from_text(content)
    if rows:
        return f'<div class="body-card table-card">{render_table(rows)}</div>'
    lines = [line.strip(" \t-•*") for line in content.splitlines() if line.strip(" \t-•*")]
    if should_render_as_list(lines):
        items = "\n".join(
            f'<li><span class="point-icon" aria-hidden="true"></span><span>{render_inline_text(line)}</span></li>'
            for line_index, line in enumerate(lines)
        )
        return f'<div class="body-card"><ul class="content-list emoji-list guided-list">{items}</ul></div>'
    return f'<div class="body-card"><p><span class="point-icon" aria-hidden="true"></span><span>{render_inline_text(content)}</span></p></div>'


def should_render_timeline(slide: Dict[str, Any], title_text: str, body_items: List[Dict[str, Any]]) -> bool:
    hints = slide.get("layoutHints") or []
    if "timeline" in hints:
        return True
    if len(body_items) < 3 or len(body_items) > 5:
        return False
    title_lower = title_text.lower()
    return any(token in title_lower for token in ("process", "step", "roadmap", "流程", "步骤", "路径", "时间"))


def is_chapter_slide(title_text: str, texts: List[Dict[str, Any]]) -> bool:
    combined = "\n".join(str(item.get("content") or "") for item in texts if item.get("kind") != "table")
    return bool(re.search(r"\bchapter\s+[0-9ivxlcdm]+\b", title_text + "\n" + combined, flags=re.IGNORECASE))


def render_timeline(body_items: List[Dict[str, Any]]) -> str:
    steps: List[str] = []
    for item in body_items[:5]:
        content = str(item.get("content") or "").strip()
        if not content:
            continue
        first_line = next((line.strip(" \t-•*") for line in content.splitlines() if line.strip(" \t-•*")), content)
        steps.append(first_line)
    if not steps:
        return ""
    return '<div class="timeline">' + "\n".join(
        f'<div class="timeline-step"><span class="step-number">{index}</span><span>{render_inline_text(step)}</span></div>'
        for index, step in enumerate(steps, start=1)
    ) + "</div>"


def render_slide(slide: Dict[str, Any], preset: Dict[str, str]) -> str:
    texts = slide["texts"]
    images = slide.get("images", [])
    title_item = next((item for item in texts if item.get("kind") != "table"), None)
    title_text = clean_title_text(title_item["content"] if title_item else "Slide")
    body_items = [item for item in texts if item is not title_item]
    layout_classes = determine_slide_classes(slide, texts, images)

    if should_render_timeline(slide, title_text, body_items):
        body_html = render_timeline(body_items)
    else:
        body_html = "\n".join(
            render_body_item(item, index) for index, item in enumerate(body_items) if item.get("content")
        )
    images_html = "\n".join(
        f"""<figure class="media-frame"><img src="{html.escape(image['src'])}" alt="{html.escape(image['alt'])}"></figure>"""
        for image in images
    )
    media_count_class = "media-count-1"
    if len(images) == 2:
        media_count_class = "media-count-2"
    elif len(images) > 2:
        media_count_class = "media-count-many"
    slide_class = " ".join(layout_classes)
    media_column = ""
    if images_html:
        media_column = f"""
    <div class="media-column {media_count_class}">
      {images_html}
    </div>"""

    return f"""<article class="{slide_class}" data-slide="{slide['page']}" data-layout="{html.escape(' '.join(layout_classes[1:]))}">
  <div class="slide-inner">
    <div class="text-column">
      <h1 class="slide-title" data-editable="true">{html.escape(title_text)}</h1>
      <div class="slide-body" data-editable="true">
        {body_html}
      </div>
    </div>{media_column}
    <div class="visual-anchor" aria-hidden="true"></div>
    <div class="page-number">{slide['page']}</div>
  </div>
</article>"""


def determine_slide_classes(
    slide: Dict[str, Any],
    texts: List[Dict[str, Any]],
    images: List[Dict[str, Any]],
) -> List[str]:
    classes = ["slide", "has-media" if images else "text-only"]
    if slide.get("page") == 1:
        classes.append("cover")
    title_item = next((item for item in texts if item.get("kind") != "table"), None)
    title_text = clean_title_text(str(title_item.get("content") or "")) if title_item else ""
    if is_chapter_slide(title_text, texts):
        classes.append("chapter")
    if len(title_text) >= 96:
        classes.append("title-very-long")
    elif len(title_text) >= 58:
        classes.append("title-long")
    generated_kind = slide.get("generatedKind")
    if generated_kind in {"agenda", "transition"}:
        classes.append(generated_kind)

    text_length = sum(len(str(item.get("content") or "")) for item in texts)
    if text_length <= 220 and len(texts) <= 3:
        classes.append("density-short")
    elif text_length >= 760 or len(texts) >= 7:
        classes.append("density-dense")
    else:
        classes.append("density-medium")

    if not images:
        apply_layout_hints(classes, slide.get("layoutHints") or [])
        if "cards" not in classes:
            classes.append("cards")
        return classes

    if len(images) >= 3:
        classes.append("gallery")
    if text_length >= 540 or len(texts) >= 5:
        classes.append("text-heavy")
    if len(images) == 1 and text_length <= 340:
        classes.append("image-dominant")

    image_side = infer_image_side(slide, images)
    classes.append(image_side)
    apply_layout_hints(classes, slide.get("layoutHints") or [])
    return classes


def apply_layout_hints(classes: List[str], hints: List[str]) -> None:
    if not hints:
        return
    density_tokens = {"density-short", "density-medium", "density-dense"}
    side_tokens = {"image-left", "image-right"}
    if any(hint in density_tokens for hint in hints):
        classes[:] = [item for item in classes if item not in density_tokens]
    if any(hint in side_tokens for hint in hints):
        classes[:] = [item for item in classes if item not in side_tokens]
    for hint in hints:
        if hint not in classes:
            classes.append(hint)


def infer_image_side(slide: Dict[str, Any], images: List[Dict[str, Any]]) -> str:
    source_width = int((slide.get("sourceSize") or {}).get("width") or 0)
    if not source_width:
        return "image-right"

    total_area = 0
    weighted_center = 0.0
    for image in images:
        width = max(1, int(image.get("width") or 1))
        height = max(1, int(image.get("height") or 1))
        area = width * height
        center_x = int(image.get("left") or 0) + width / 2
        weighted_center += center_x * area
        total_area += area

    if not total_area:
        return "image-right"
    ratio = weighted_center / total_area / source_width
    if ratio < 0.46:
        return "image-left"
    return "image-right"


def load_jobs(jobs_path: Path) -> List[Dict[str, Any]]:
    if not jobs_path.exists():
        return []
    try:
        return json.loads(jobs_path.read_text(encoding="utf-8-sig"))
    except Exception:
        backup = jobs_path.with_suffix(f".broken-{int(time.time())}.json")
        shutil.copy2(jobs_path, backup)
        return []


def save_jobs(jobs_path: Path, jobs: List[Dict[str, Any]]) -> None:
    jobs_path.parent.mkdir(parents=True, exist_ok=True)
    jobs_path.write_text(json.dumps(jobs, ensure_ascii=False, indent=2), encoding="utf-8-sig")

